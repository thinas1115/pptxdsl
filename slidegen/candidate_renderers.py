"""採用判断用に追加した業務スライドrenderer群。"""

from collections import defaultdict
from itertools import permutations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from asset_paths import resolve_icon_path
from diagrams import add_arrow
from diagrams3 import plain_line, route
from generate import (
    ACCENT,
    BODY_W,
    CANVAS,
    GRAY,
    LIGHT,
    MARGIN,
    NAVY,
    RULE,
    TEXT,
    WHITE,
    ZEBRA,
    add_rect,
    add_text,
    header,
)
from layout_fit import FitError, fit_text_or_raise, select_fit, stepped
from quality_markers import SURFACE_ON_CANVAS_PREFIX
from textfit import text_width_in, wrap_natural


def _flatten_shape(shape):
    """テーマ由来の影・光彩を外し、環境差のないフラットな図形にする。"""
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)
    shape.shadow.inherit = False
    return shape


def _flat_rect(slide, x, y, w, h, fill, *, line=None, round_=False):
    return _flatten_shape(
        add_rect(slide, x, y, w, h, fill, line=line, round_=round_))


def _fit_rows(renderer, available, count, *, row_h, min_row_h, gap, min_gap,
              font, min_font, reserve=0.0):
    """余白、行高、文字の順に縮める共通収容処理。"""
    def used(values):
        return reserve + count * values["row_h"] + max(0, count - 1) * values["gap"]

    def candidates():
        yield "standard", {"row_h": row_h, "gap": gap, "font": font}, used(
            {"row_h": row_h, "gap": gap, "font": font})
        for current_gap in stepped(gap - 0.03, min_gap, 0.03):
            values = {"row_h": row_h, "gap": current_gap, "font": font}
            yield "gap", values, used(values)
        for current_h in stepped(row_h - 0.04, min_row_h, 0.04):
            values = {"row_h": current_h, "gap": min_gap, "font": font}
            yield "element", values, used(values)
        for current_font in stepped(font - 0.5, min_font, 0.5):
            values = {"row_h": min_row_h, "gap": min_gap, "font": current_font}
            yield "font", values, used(values)

    return select_fit(
        renderer,
        available,
        candidates(),
        guidance="項目を減らすか、文言を短くするか、複数スライドへ分割してください。",
    )


def _text_in_box(slide, renderer, field, x, y, w, h, text, max_pt, min_pt,
                 *, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, spacing=1.15, role=None):
    size, lines = fit_text_or_raise(
        renderer,
        field,
        text,
        w,
        h,
        max_pt,
        min_pt=min_pt,
        weight="bold" if bold else "regular",
        spacing=spacing,
        role=role,
    )
    return add_text(
        slide, x, y, w, h, "\n".join(lines), size,
        bold=bold, color=color, align=align, anchor=anchor, spacing=spacing,
    )


def _dot(slide, cx, cy, color=ACCENT, size=0.10):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - size / 2),
        Inches(cy - size / 2),
        Inches(size),
        Inches(size),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return _flatten_shape(shape)


def _icon_medallion(slide, x, y, icon, *, size=0.46, muted=False):
    """Fluentアイコンを余白込みの円へ中央配置する。"""
    disk = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    disk.fill.solid()
    disk.fill.fore_color.rgb = ZEBRA if muted else LIGHT
    disk.line.fill.background()
    _flatten_shape(disk)
    icon_path = resolve_icon_path(f"icons/fluent/{icon}.png")
    icon_size = size * 0.52
    slide.shapes.add_picture(
        str(icon_path), Inches(x + (size - icon_size) / 2),
        Inches(y + (size - icon_size) / 2), Inches(icon_size), Inches(icon_size))


def _status_medallion(slide, x, y, mark, *, positive=True, size=0.42):
    disk = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    disk.fill.solid()
    disk.fill.fore_color.rgb = ACCENT if positive else GRAY
    disk.line.fill.background()
    _flatten_shape(disk)
    add_text(slide, x, y - 0.01, size, size, mark, 16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _number_medallion(slide, x, y, number, *, size=0.46):
    disk = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    disk.fill.solid()
    disk.fill.fore_color.rgb = LIGHT
    disk.line.color.rgb = RULE
    _flatten_shape(disk)
    add_text(slide, x, y, size, size, f"{number:02d}", 9.5,
             bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def s_scope(slide, spec, page):
    """対象範囲と対象外を2列で整理する。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    left_items, right_items = spec["in_scope"], spec["out_of_scope"]
    body_top = area.top + 0.20
    body_bottom = area.bottom - 0.14
    max_count = max(len(left_items), len(right_items))
    sparse = max_count <= 2
    row_h = 1.04 if sparse else 0.58
    row_gap = 0.16 if sparse else 0.07
    font = 17.0 if sparse else 14.0
    fitted = _fit_rows(
        "scope", body_bottom - body_top - 0.82, max_count,
        row_h=row_h, min_row_h=0.44, gap=row_gap, min_gap=0.03,
        font=font, min_font=11.0,
    )
    values = fitted.values
    col_gap, col_w = 0.34, (BODY_W - 0.34) / 2
    content_h = max_count * values["row_h"] + max(0, max_count - 1) * values["gap"]
    panel_h = 0.70 + content_h
    if sparse:
        panel_h = max(panel_h, 2.22)
        body_top += max(0.0, body_bottom - body_top - panel_h) * 0.26

    def column(x, label, items, positive):
        color = ACCENT if positive else GRAY
        # 見出しだけが白く浮かないよう、外周を持つ一枚のパネルとして描く。
        _flat_rect(slide, x, body_top, col_w, panel_h, WHITE, line=RULE)
        _flat_rect(slide, x, body_top, col_w, 0.62,
                   LIGHT if positive else ZEBRA)
        _status_medallion(slide, x + 0.20, body_top + 0.10,
                          "✓" if positive else "×", positive=positive)
        add_text(slide, x + 0.76, body_top + 0.16, col_w - 0.98, 0.28, label,
                 15.0, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        plain_line(slide, x, body_top + 0.62, x + col_w, body_top + 0.62,
                   color=RULE, width=0.8)
        y = body_top + 0.70 + max(0.0, panel_h - 0.70 - content_h) / 2
        for index, item in enumerate(items):
            marker = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + 0.25), Inches(y + values["row_h"] / 2 - 0.035),
                Inches(0.07), Inches(0.07))
            marker.fill.solid()
            marker.fill.fore_color.rgb = color
            marker.line.fill.background()
            _flatten_shape(marker)
            _text_in_box(
                slide, "scope", f"items[{index}]", x + 0.48, y,
                col_w - 0.68, values["row_h"], item,
                values["font"], 10.5, color=TEXT,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            if index < len(items) - 1:
                plain_line(slide, x + 0.48, y + values["row_h"],
                           x + col_w - 0.24, y + values["row_h"],
                           color=RULE, width=0.45)
            y += values["row_h"] + values["gap"]

    column(MARGIN, spec.get("in_label", "対象"), left_items, True)
    column(MARGIN + col_w + col_gap,
           spec.get("out_label", "対象外"), right_items, False)


def s_summary(slide, spec, page):
    """少数の論点を、視線順が明確な要約面へまとめる。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    sections = spec["sections"]
    top = area.top + 0.26
    available = area.bottom - top - 0.22
    count = len(sections)
    if count <= 3:
        cols, rows = count, 1
    else:
        cols, rows = 2, 2
    gap_x, gap_y = 0.0, 0.0
    cell_w = (BODY_W - gap_x * (cols - 1)) / cols
    # 短い要約を本文下端まで引き延ばさず、内容の自然なまとまりを保つ。
    natural_h = (2.70 if count == 2 else 2.46) if rows == 1 else 1.78
    cell_h = min(natural_h, (available - gap_y * (rows - 1)) / rows)
    if cell_h < 1.22:
        raise FitError("summary: 本文領域へ論点を配置できません。論点を減らしてください。")

    block_h = rows * cell_h + (rows - 1) * gap_y
    extra = max(0.0, area.bottom - top - block_h)
    top += extra * 0.34

    heading_pt = 18.5 if count == 2 else (16.5 if count == 3 else 15.5)
    body_pt = 15.0 if count == 2 else (13.0 if count == 3 else 12.0)
    if rows == 2:
        _flat_rect(slide, MARGIN, top, BODY_W, rows * cell_h, WHITE)
        plain_line(slide, MARGIN + BODY_W / 2, top + 0.18,
                   MARGIN + BODY_W / 2, top + rows * cell_h - 0.18,
                   color=RULE, width=0.65)
        plain_line(slide, MARGIN + 0.18, top + cell_h,
                   MARGIN + BODY_W - 0.18, top + cell_h,
                   color=RULE, width=0.65)
    for index, section in enumerate(sections):
        row, col = divmod(index, cols)
        x = MARGIN + col * (cell_w + gap_x)
        y = top + row * (cell_h + gap_y)
        if rows == 1 and index > 0:
            plain_line(slide, x, y + 0.18, x, y + cell_h - 0.18,
                       color=RULE, width=0.65)
        if section.get("icon"):
            _icon_medallion(slide, x + 0.24, y + 0.20,
                            section["icon"], size=0.50)
        else:
            _number_medallion(slide, x + 0.24, y + 0.22,
                              index + 1, size=0.46)
        _text_in_box(slide, "summary", f"sections[{index}].heading",
                     x + 0.92, y + 0.24, cell_w - 1.16, 0.32,
                     section["heading"], heading_pt, 12.0,
                     bold=True, color=NAVY)
        plain_line(slide, x + 0.24, y + 0.84, x + cell_w - 0.24, y + 0.84,
                   color=ACCENT if section.get("icon") else RULE, width=0.8)
        _text_in_box(slide, "summary", f"sections[{index}].body",
                     x + 0.26, y + 1.00, cell_w - 0.52, cell_h - 1.18,
                     section["body"], body_pt, 10.0, color=TEXT, spacing=1.22)
def s_paired_comparison(slide, spec, page):
    """同一観点で左右を比較し、対応関係を行単位で固定する。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    rows = spec["rows"]
    top = area.top + 0.20
    sparse = len(rows) <= 2
    row_h = 0.88 if sparse else 0.68
    row_gap = 0.0
    font = 15.0 if sparse else 13.5
    fitted = _fit_rows(
        "paired_comparison", area.bottom - top - 0.54,
        len(rows), row_h=row_h, min_row_h=0.46, gap=row_gap, min_gap=0.0,
        font=font, min_font=10.0,
    )
    values = fitted.values
    rows_h = len(rows) * values["row_h"] + max(0, len(rows) - 1) * values["gap"]
    block_h = 0.70 + rows_h
    top += max(0.0, area.bottom - top - block_h) * 0.34
    left_x, left_w = MARGIN, 4.95
    criterion_x, criterion_w = left_x + left_w + 0.20, 1.68
    right_x = criterion_x + criterion_w + 0.20
    right_w = MARGIN + BODY_W - right_x
    _flat_rect(slide, left_x, top, left_w, 0.56, NAVY, round_=True)
    _flat_rect(slide, right_x, top, right_w, 0.56, ACCENT, round_=True)
    add_text(slide, left_x + 0.18, top, left_w - 0.36, 0.56,
             spec["left_label"], 15.5, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, right_x + 0.18, top, right_w - 0.36, 0.56,
             spec["right_label"], 15.5, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, criterion_x, top + 0.13, criterion_w, 0.28,
             spec.get("criterion_label", "評価軸"), 10.0,
             bold=True, color=GRAY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    _flat_rect(slide, left_x, top + 0.64, left_w, rows_h, ZEBRA)
    _flat_rect(slide, right_x, top + 0.64, right_w, rows_h, LIGHT)
    y = top + 0.64
    for index, row in enumerate(rows):
        if index > 0:
            plain_line(slide, left_x + 0.18, y,
                       left_x + left_w - 0.18, y, color=RULE, width=0.55)
            plain_line(slide, right_x + 0.18, y,
                       right_x + right_w - 0.18, y, color=RULE, width=0.55)
        cy = y + values["row_h"] / 2
        _dot(slide, left_x + left_w, cy, GRAY, 0.08)
        _dot(slide, right_x, cy, ACCENT, 0.08)
        plain_line(slide, left_x + left_w, cy, criterion_x, cy,
                   color=RULE, width=0.8)
        plain_line(slide, criterion_x + criterion_w, cy, right_x, cy,
                   color=RULE, width=0.8)
        _text_in_box(slide, "paired_comparison", f"rows[{index}].left",
                     left_x + 0.20, y, left_w - 0.40, values["row_h"], row["left"],
                     values["font"], 9.5, anchor=MSO_ANCHOR.MIDDLE)
        _flat_rect(slide, criterion_x + 0.06, y + 0.12,
                   criterion_w - 0.12, values["row_h"] - 0.24, WHITE)
        _text_in_box(slide, "paired_comparison", f"rows[{index}].criterion",
                     criterion_x + 0.06, y + 0.12,
                     criterion_w - 0.12, values["row_h"] - 0.24, row["criterion"],
                     min(11.0, values["font"]), 8.5, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     role="compact")
        _text_in_box(slide, "paired_comparison", f"rows[{index}].right",
                     right_x + 0.20, y, right_w - 0.40, values["row_h"], row["right"],
                     values["font"], 9.5, anchor=MSO_ANCHOR.MIDDLE)
        y += values["row_h"] + values["gap"]
def _mapping_items_by_min_crossings(left_items, right_items, links):
    """関係を変えずに左右項目を並べ替え、直接結線の交差数を最小化する。"""
    left_ids = [item["id"] for item in left_items]
    right_ids = [item["id"] for item in right_items]
    original_left = {item_id: index for index, item_id in enumerate(left_ids)}
    original_right = {item_id: index for index, item_id in enumerate(right_ids)}
    pairs = [(link["from"], link["to"]) for link in links]
    best = None
    for left_order in permutations(left_ids):
        left_rank = {item_id: index for index, item_id in enumerate(left_order)}
        for right_order in permutations(right_ids):
            right_rank = {item_id: index for index, item_id in enumerate(right_order)}
            crossings = sum(
                (left_rank[source_a] - left_rank[source_b])
                * (right_rank[target_a] - right_rank[target_b]) < 0
                for index, (source_a, target_a) in enumerate(pairs)
                for source_b, target_b in pairs[index + 1:]
            )
            displacement = sum(
                abs(index - original_left[item_id])
                for index, item_id in enumerate(left_order)
            ) + sum(
                abs(index - original_right[item_id])
                for index, item_id in enumerate(right_order)
            )
            score = (crossings, displacement, left_order, right_order)
            if best is None or score < best:
                best = score
                if crossings == 0 and displacement == 0:
                    break
        if best[:2] == (0, 0):
            break
    left_by_id = {item["id"]: item for item in left_items}
    right_by_id = {item["id"]: item for item in right_items}
    return (
        [left_by_id[item_id] for item_id in best[2]],
        [right_by_id[item_id] for item_id in best[3]],
        best[0],
    )


def s_mapping(slide, spec, page):
    """左右の項目間に一対多・多対多の対応線を描く。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    left_items, right_items = spec["left_items"], spec["right_items"]
    links = spec["links"]
    left_items, right_items, _crossings = _mapping_items_by_min_crossings(
        left_items, right_items, links)
    top = area.top + 0.22
    body_bottom = area.bottom - 0.10
    max_count = max(len(left_items), len(right_items))
    sparse = max_count <= 2
    fitted = _fit_rows(
        "mapping", body_bottom - top - 0.52, max_count,
        row_h=0.82 if sparse else 0.58,
        min_row_h=0.58 if sparse else 0.42,
        gap=0.22 if sparse else 0.13,
        min_gap=0.12 if sparse else 0.04,
        font=15.0 if sparse else 13.5, min_font=10.0,
    )
    values = fitted.values
    left_x, left_w = MARGIN + 0.16, 4.60
    right_x = MARGIN + BODY_W - left_w - 0.16
    right_w = left_w
    header_h = 0.48
    for x, width, label, text_x in (
        (left_x, left_w, spec["left_label"], left_x + 0.52),
        (right_x, right_w, spec["right_label"], right_x + 0.18),
    ):
        add_text(slide, text_x, top, width - 0.70, header_h, label, 14.0,
                 bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        plain_line(slide, x, top + header_h, x + width, top + header_h,
                   color=ACCENT, width=1.1)
    start_y = top + header_h + 0.08

    def positions(items):
        return {
            item["id"]: start_y + i * (values["row_h"] + values["gap"])
            for i, item in enumerate(items)
        }

    left_y, right_y = positions(left_items), positions(right_items)
    def draw_items(items, y_by_id, x, width, side):
        for index, item in enumerate(items):
            y = y_by_id[item["id"]]
            _flat_rect(slide, x, y, width, values["row_h"],
                       ZEBRA if side == "left" else LIGHT)
            marker_x = x + 0.25 if side == "left" else x + width - 0.25
            _dot(slide, marker_x, y + values["row_h"] / 2, ACCENT, 0.25)
            add_text(slide, marker_x - 0.12, y + values["row_h"] / 2 - 0.12,
                     0.24, 0.24, f"{index + 1}", 8.5, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            text_x = x + 0.52 if side == "left" else x + 0.18
            text_w = width - 0.70
            _text_in_box(slide, "mapping", f"{side}_items[{index}].text",
                         text_x, y, text_w, values["row_h"], item["text"],
                         values["font"], 9.5, bold=True, color=NAVY,
                         anchor=MSO_ANCHOR.MIDDLE)
            if index < len(items) - 1:
                plain_line(slide, text_x, y + values["row_h"],
                           x + width - 0.18 if side == "left" else marker_x - 0.34,
                           y + values["row_h"], color=RULE, width=0.45)

    gap_left = left_x + left_w
    gap_right = right_x
    endpoint_colors = {}
    # 参考スライドと同じく、項目間を直接追えることを優先する。線を先に描き、
    # 項目と端点を前面へ置くことで、交差があっても接続元と接続先を見失いにくくする。
    for link in links:
        sy = left_y[link["from"]] + values["row_h"] / 2
        ty = right_y[link["to"]] + values["row_h"] / 2
        color = ACCENT if link.get("emphasis") else GRAY
        width = 1.4 if link.get("emphasis") else 1.0
        # 後から描く線の下へ背景色の縁を置き、交点で線が連結して見える
        # 誤読を防ぐ。端点は項目描画後に置き直す。
        plain_line(slide, gap_left, sy, gap_right, ty,
                   color=CANVAS, width=3.6)
        add_arrow(slide, gap_left, sy, gap_right, ty,
                  color=color, width=width)
        endpoint_colors[("left", link["from"])] = color
        endpoint_colors[("right", link["to"])] = color
    draw_items(left_items, left_y, left_x, left_w, "left")
    draw_items(right_items, right_y, right_x, right_w, "right")
    for item in left_items:
        color = endpoint_colors.get(("left", item["id"]), GRAY)
        _dot(slide, gap_left, left_y[item["id"]] + values["row_h"] / 2,
             color, 0.08)
    for item in right_items:
        color = endpoint_colors.get(("right", item["id"]), GRAY)
        _dot(slide, gap_right, right_y[item["id"]] + values["row_h"] / 2,
             color, 0.08)
def _swimlane_step_rects(spec, x0, y0, stage_w, lane_h):
    grouped = defaultdict(list)
    lane_index = {lane["id"]: index for index, lane in enumerate(spec["lanes"])}
    stage_index = {stage["id"]: index for index, stage in enumerate(spec["stages"])}
    for step in spec["steps"]:
        grouped[(step["lane"], step["stage"])].append(step)
    rects = {}
    for (lane_id, stage_id), steps in grouped.items():
        cell_x = x0 + stage_index[stage_id] * stage_w
        cell_y = y0 + lane_index[lane_id] * lane_h
        gap = 0.10
        if len(spec["stages"]) <= 2:
            max_node_w = 2.36
        elif len(spec["stages"]) <= 3:
            max_node_w = 1.92
        else:
            max_node_w = 1.55
        available_w = stage_w - 0.28 - gap * (len(steps) - 1)
        widths = [
            min(max_node_w, max(
                0.88, text_width_in(step["name"], 10.5, "bold") + 0.36))
            for step in steps
        ]
        if sum(widths) > available_w:
            scale = available_w / sum(widths)
            widths = [width * scale for width in widths]
        total_w = sum(widths) + gap * (len(steps) - 1)
        start_x = cell_x + (stage_w - total_w) / 2
        node_h = min(0.62 if lane_h >= 1.20 else 0.50, lane_h - 0.24)
        cursor_x = start_x
        for step, width in zip(steps, widths):
            rects[step["id"]] = (
                cursor_x,
                cell_y + (lane_h - node_h) / 2,
                width,
                node_h,
            )
            cursor_x += width + gap
    return rects


def _circled_step_number(index):
    """工程順を視線の入口にする。20件超は通常数字へ安全にフォールバックする。"""
    return chr(0x2460 + index) if 0 <= index < 20 else f"{index + 1}."


def _swimlane_step_label(step, index):
    number = step.get("number", index + 1)
    if number is None:
        return step["name"]
    return f"{_circled_step_number(number - 1)} {step['name']}"


def _swimlane_header(slide, spec):
    """参照デザインの縦リズムを保つswimlane専用ヘッダー。"""
    _flat_rect(slide, 0, 0, 13.333, 7.5, CANVAS)
    plain_line(slide, 0.19, 0.20, 13.08, 0.20,
               color=ACCENT, width=3.0)
    title_size, title_lines = fit_text_or_raise(
        "swimlane", "title", spec["title"], 12.25, 0.48, 25.5,
        min_pt=19.0, weight="bold", spacing=1.08,
    )
    add_text(slide, 0.53, 0.54, 12.25, 0.48, "\n".join(title_lines),
             title_size, bold=True, color=NAVY, spacing=1.08)
    if spec.get("lead"):
        lead_size, lead_lines = fit_text_or_raise(
            "swimlane", "lead", spec["lead"], 12.25, 0.34, 14.0,
            min_pt=11.5, spacing=1.12,
        )
        add_text(slide, 0.53, 1.05, 12.25, 0.34, "\n".join(lead_lines),
                 lead_size, color=GRAY, spacing=1.12)


def s_swimlane(slide, spec, page):
    """担当レーンと工程フェーズを持つ業務フローを描く。"""
    _swimlane_header(slide, spec)
    lanes, stages = spec["lanes"], spec["stages"]
    frame_x, frame_w = 0.37, 12.59
    top = 1.64
    stage_h, lane_label_w = 0.37, 1.27
    stage_index = {stage["id"]: index for index, stage in enumerate(stages)}
    step_by_id = {step["id"]: step for step in spec["steps"]}
    has_feedback = any(
        stage_index[step_by_id[edge["to"]]["stage"]]
        < stage_index[step_by_id[edge["from"]]["stage"]]
        for edge in spec["edges"]
    )
    edge_kinds = {
        "feedback" if edge.get("kind") == "feedback" else "forward"
        for edge in spec["edges"]
    }
    feedback_h = 0.30 if has_feedback else 0.0
    show_line_legend = len(edge_kinds) > 1
    body_bottom = 6.48
    available_h = body_bottom - top - stage_h
    standard_lane_h = (
        1.34 if len(lanes) == 2 else 0.98 if len(lanes) == 3 else 0.95)
    fitted = _fit_rows(
        "swimlane", available_h, len(lanes),
        row_h=standard_lane_h, min_row_h=0.62, gap=0.0, min_gap=0.0,
        font=11.5, min_font=9.0, reserve=feedback_h,
    )
    lane_h = fitted.values["row_h"]
    x0 = frame_x + lane_label_w
    stage_w = (frame_w - lane_label_w) / len(stages)
    lane_y0 = top + stage_h
    for index, stage in enumerate(stages):
        x = x0 + index * stage_w
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, Inches(x), Inches(top),
            Inches(stage_w + (0.04 if index < len(stages) - 1 else 0)),
            Inches(stage_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RULE
        shape.line.width = Pt(0.85)
        shape.name = (
            f"{SURFACE_ON_CANVAS_PREFIX}swimlane-stage:{stage['id']}")
        _flatten_shape(shape)
        add_text(slide, x + 0.14, top + 0.06, stage_w - 0.28, 0.26,
                 stage["label"], 10.5, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if show_line_legend:
        legend_x = frame_x + 0.08
        add_text(slide, legend_x, top + 0.01, 0.28, 0.18,
                 "凡例", 7.5, bold=True, color=GRAY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_arrow(slide, legend_x + 0.36, top + 0.13,
                  legend_x + 0.62, top + 0.13, color=ACCENT, width=1.15)
        add_text(slide, legend_x + 0.68, top + 0.04, 0.42, 0.18,
                 "順方向", 7.0, color=TEXT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_arrow(slide, legend_x + 0.36, top + 0.31,
                  legend_x + 0.62, top + 0.31, color=ACCENT, width=1.15,
                  dash="dash")
        add_text(slide, legend_x + 0.68, top + 0.22, 0.42, 0.18,
                 "差戻し", 7.0, color=TEXT,
                 anchor=MSO_ANCHOR.MIDDLE)
    rects = _swimlane_step_rects(spec, x0, lane_y0, stage_w, lane_h)
    lane_index = {lane["id"]: index for index, lane in enumerate(lanes)}
    for index, lane in enumerate(lanes):
        y = lane_y0 + index * lane_h
        _flat_rect(slide, frame_x, y, lane_label_w, lane_h, LIGHT)
        _flat_rect(slide, x0, y, frame_w - lane_label_w, lane_h, WHITE)
        lane_text_w = lane_label_w - 0.22
        lane_font = 13.0
        while (lane_font > 8.0
               and text_width_in(lane["label"], lane_font, "bold") > lane_text_w):
            lane_font -= 0.5
        if text_width_in(lane["label"], lane_font, "bold") <= lane_text_w:
            lane_lines = [lane["label"]]
        else:
            lane_font = 12.0
            lane_lines = wrap_natural(
                lane["label"], lane_text_w, lane_font, "bold")
            while len(lane_lines) > 2 and lane_font > 8.0:
                lane_font -= 0.5
                lane_lines = wrap_natural(
                    lane["label"], lane_text_w, lane_font, "bold")
        if len(lane_lines) > 2:
            raise FitError(
                f"swimlane: lanes[{index}].label を2行以内へ収容できません。"
                "担当名を短くしてください。")
        add_text(slide, frame_x + 0.12, y, lane_text_w, lane_h,
                 "\n".join(lane_lines), lane_font, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        plain_line(slide, frame_x, y + lane_h, frame_x + frame_w, y + lane_h,
                   color=RULE, width=0.65)
    plain_line(slide, x0, lane_y0, x0, lane_y0 + len(lanes) * lane_h,
               color=RULE, width=0.8)
    for index in range(1, len(stages)):
        x = x0 + index * stage_w
        plain_line(slide, x, lane_y0, x, lane_y0 + len(lanes) * lane_h,
                   color=RULE, width=0.55)
    # 接続線を先に描き、ノードを上から重ねて線の貫通を隠す。
    for edge in spec["edges"]:
        sx, sy, sw, sh = rects[edge["from"]]
        tx, ty, tw, th = rects[edge["to"]]
        source_stage = stage_index[step_by_id[edge["from"]]["stage"]]
        target_stage = stage_index[step_by_id[edge["to"]]["stage"]]
        if target_stage > source_stage:
            start, end = (sx + sw, sy + sh / 2), (tx, ty + th / 2)
            if abs(start[1] - end[1]) < 0.02:
                points = [start, end]
            else:
                boundary = x0 + target_stage * stage_w - 0.04
                points = [start, (boundary, start[1]), (boundary, end[1]), end]
        elif target_stage == source_stage:
            source_lane = lane_index[step_by_id[edge["from"]]["lane"]]
            target_lane = lane_index[step_by_id[edge["to"]]["lane"]]
            if source_lane == target_lane:
                left_to_right = tx > sx
                start = (sx + sw if left_to_right else sx, sy + sh / 2)
                end = (tx if left_to_right else tx + tw, ty + th / 2)
                points = [start, end]
            elif abs(target_lane - source_lane) == 1:
                downward = ty > sy
                start = (sx + sw / 2, sy + sh if downward else sy)
                end = (tx + tw / 2, ty if downward else ty + th)
                channel_y = (start[1] + end[1]) / 2
                points = [start, (start[0], channel_y),
                          (end[0], channel_y), end]
            else:
                # 中間レーンのノードを貫通しないよう工程セル右端を通し、
                # 対象ノードの上下辺へ最後の短い縦線で接続する。
                upward = target_lane < source_lane
                channel_x = x0 + (source_stage + 1) * stage_w - 0.22
                start = (sx + sw, sy + sh / 2)
                target_x = tx + tw / 2
                target_y = ty + th if upward else ty
                approach_y = target_y + (0.10 if upward else -0.10)
                points = [start, (channel_x, start[1]),
                          (channel_x, approach_y),
                          (target_x, approach_y), (target_x, target_y)]
        else:
            channel_y = lane_y0 + len(lanes) * lane_h + 0.10
            if channel_y > body_bottom - 0.08:
                raise FitError("swimlane: 差戻し線の配線領域が不足しています。工程を分割してください。")
            start = (sx + sw / 2, sy + sh)
            end = (tx + tw / 2, ty + th)
            points = [start, (start[0], channel_y), (end[0], channel_y), end]
        route(slide, points,
              dash="dash" if edge.get("kind") == "feedback" else None,
              width=1.10, color=ACCENT)
    for index, step in enumerate(spec["steps"]):
        x, y, w, h = rects[step["id"]]
        _flat_rect(slide, x, y, w, h,
                   LIGHT if step.get("style") == "accent" else WHITE,
                   line=ACCENT)
        display_name = _swimlane_step_label(step, index)
        one_line_size = min(11.5, fitted.values["font"])
        text_x = x + 0.06
        text_w = w - 0.12
        while (one_line_size >= 7.5
               and text_width_in(display_name, one_line_size, "bold") > text_w):
            one_line_size -= 0.5
        if one_line_size >= 7.5:
            add_text(slide, text_x, y, text_w, h, display_name,
                     one_line_size, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            _text_in_box(slide, "swimlane", f"steps[{index}].name",
                         text_x, y, text_w, h, display_name,
                         fitted.values["font"], 7.5, bold=True, color=NAVY,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                         role="compact")
def s_sequence(slide, spec, page):
    """参加者間のメッセージを上から下へ時系列表示する。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    participants, messages = spec["participants"], spec["messages"]
    phases = spec.get("phases", [])
    kinds = {message.get("kind", "request") for message in messages}
    show_legend = len(kinds) > 1
    legend_h = 0.16 if show_legend else 0.0
    top = area.top + 0.16
    header_h = 0.46
    message_top = top + header_h + 0.14
    available = area.bottom - message_top - legend_h - 0.12
    sparse = len(messages) <= 4
    fitted = _fit_rows(
        "sequence", available, len(messages),
        row_h=0.78 if sparse else 0.42,
        min_row_h=0.56 if sparse else 0.29,
        gap=0.34 if sparse else 0.12,
        min_gap=0.16 if sparse else 0.02,
        font=13.5 if sparse else 10.5, min_font=8.0,
    )
    sequence_center = MARGIN + BODY_W / 2
    half_span = (4.18 if len(participants) <= 3
                 else min(4.45, 2.80 + 0.55 * (len(participants) - 2)))
    participant_x0 = sequence_center - half_span
    participant_x1 = sequence_center + half_span
    used_h = (len(messages) * fitted.values["row_h"]
              + max(0, len(messages) - 1) * fitted.values["gap"])
    message_top += max(0.0, available - used_h) * 0.28
    xs = [participant_x0 + i * (participant_x1 - participant_x0) / (len(participants) - 1)
          for i in range(len(participants))]
    x_by_id = {participant["id"]: x for participant, x in zip(participants, xs)}
    message_y = {
        message["id"]: message_top + i * (fitted.values["row_h"] + fitted.values["gap"])
        for i, message in enumerate(messages)
    }
    index_by_id = {message["id"]: index for index, message in enumerate(messages)}
    message_fill = {message["id"]: CANVAS for message in messages}
    lifeline_bottom = message_y[messages[-1]["id"]] + fitted.values["row_h"] + 0.08
    for phase_index, phase in enumerate(phases):
        start_index = index_by_id[phase["from"]]
        end_index = index_by_id[phase["to"]]
        start_center = (message_y[messages[start_index]["id"]]
                        + fitted.values["row_h"] / 2)
        # ラベル上端をフェーズ境界にし、直前フェーズの背景が文字へ重ならないようにする。
        y1 = start_center - 0.28
        if end_index + 1 < len(messages):
            next_center = (message_y[messages[end_index + 1]["id"]]
                           + fitted.values["row_h"] / 2)
            y2 = next_center - 0.29
        else:
            y2 = lifeline_bottom
        phase_x = participant_x0 - 0.48
        phase_w = participant_x1 - participant_x0 + 0.96
        phase_fill = LIGHT if phase_index % 2 == 0 else ZEBRA
        _flat_rect(slide, phase_x, y1, phase_w, y2 - y1,
                   phase_fill)
        for message in messages[start_index:end_index + 1]:
            message_fill[message["id"]] = phase_fill
        phase_box = add_text(
            slide, phase_x + 0.14, y1 + 0.02, 1.12, 0.22,
            phase["label"], 8.5, bold=True, color=ACCENT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = LIGHT if phase_index % 2 == 0 else ZEBRA
    for index, (participant, x) in enumerate(zip(participants, xs)):
        box_w = min(2.08 if sparse else 1.74,
                    (participant_x1 - participant_x0) / len(participants) * 0.88)
        _text_in_box(slide, "sequence", f"participants[{index}].label",
                     x - box_w / 2 + 0.08, top, box_w - 0.16, header_h,
                     participant["label"], 12.5 if sparse else 11.0, 8.5,
                     bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     role="compact")
        plain_line(slide, x - box_w / 2 + 0.08, top + header_h,
                   x + box_w / 2 - 0.08, top + header_h,
                   color=RULE, width=0.8)
        plain_line(slide, x, top + header_h, x, lifeline_bottom,
                   color=GRAY, width=0.75, dash="dash")
    for index, message in enumerate(messages):
        y = message_y[message["id"]] + fitted.values["row_h"] / 2
        sx, tx = x_by_id[message["from"]], x_by_id[message["to"]]
        if sx == tx:
            loop_w, loop_h = 0.46, max(0.22, fitted.values["row_h"] * 0.72)
            points = [(sx, y - loop_h / 2), (sx + loop_w, y - loop_h / 2),
                      (sx + loop_w, y + loop_h / 2), (sx, y + loop_h / 2)]
            route(slide, points,
                  dash="dash" if message.get("kind") == "return" else None,
                  width=1.0)
            label_x, label_w = sx + 0.10, 1.20
        else:
            kind = message.get("kind", "request")
            add_arrow(slide, sx, y, tx, y,
                      color=ACCENT if kind == "async" else GRAY,
                      width=1.20 if kind == "async" else 1.05,
                      dash="dash" if kind == "return" else None)
            label_x, label_w = (sx + tx) / 2, min(1.72, abs(tx - sx) - 0.10)
        label_y = y - 0.27
        label_size, lines = fit_text_or_raise(
            "sequence", f"messages[{index}].label", message["label"],
            label_w, 0.24, fitted.values["font"], min_pt=7.5,
            weight="bold", spacing=1.0, role="compact")
        tb = add_text(slide, label_x - label_w / 2, label_y, label_w, 0.24,
                      "\n".join(lines), label_size, bold=True, color=NAVY,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                      spacing=1.0)
        tb.fill.solid()
        tb.fill.fore_color.rgb = message_fill[message["id"]]
    if show_legend:
        legend_y = area.bottom - 0.17
        legend_x = MARGIN + 0.10
        add_text(slide, legend_x, legend_y, 0.38, 0.18, "凡例", 7.5,
                 bold=True, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
        entries = [("要求", GRAY, None), ("応答", GRAY, "dash")]
        if "async" in kinds:
            entries.append(("非同期", ACCENT, None))
        for index, (label, color, dash) in enumerate(entries):
            x = legend_x + 0.48 + index * 1.05
            add_arrow(slide, x, legend_y + 0.09, x + 0.28, legend_y + 0.09,
                      color=color, width=1.0, dash=dash)
            add_text(slide, x + 0.34, legend_y, 0.56, 0.18, label, 7.5,
                     color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
