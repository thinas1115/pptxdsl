"""候補rendererの標準・過密入力と段階的な収容処理を検証する。"""

from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches

import generate
from candidate_renderers import (
    _fit_rows,
    _mapping_items_by_min_crossings,
    s_mapping,
    s_paired_comparison,
    s_scope,
    s_sequence,
    s_summary,
    s_swimlane,
)
from candidate_review_cases import REVIEW_DECK
from content_patterns import PATTERN_DECK
from layout_fit import FitError
from validate_content import validate


RENDERERS = {
    "scope": s_scope,
    "summary": s_summary,
    "paired_comparison": s_paired_comparison,
    "mapping": s_mapping,
    "swimlane": s_swimlane,
    "sequence": s_sequence,
}


def _presentation():
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    return prs


def _render(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    RENDERERS[spec["type"]](slide, spec, len(prs.slides))
    return slide


def _assert_in_slide(slide):
    tolerance = Inches(0.01)
    slide_width = Inches(generate.SLIDE_W)
    slide_height = Inches(generate.SLIDE_H)
    for shape in slide.shapes:
        assert shape.left >= -tolerance
        assert shape.top >= -tolerance
        assert shape.left + shape.width <= slide_width + tolerance
        assert shape.top + shape.height <= slide_height + tolerance


def _base(type_):
    return {"type": type_, "kicker": "検証", "title": "最大件数でも読みやすさを維持する"}


def _dense_specs():
    scope = dict(
        _base("scope"),
        in_scope=[f"実施対象{i + 1}の要件と作業範囲" for i in range(6)],
        out_of_scope=[f"対象外{i + 1}の責任範囲" for i in range(6)],
        assumptions=[f"前提条件{i + 1}" for i in range(4)],
    )
    summary = dict(
        _base("summary"),
        sections=[
            {"heading": f"論点{i + 1}", "body": "判断に必要な事実と示唆を簡潔に整理する。"}
            for i in range(4)
        ],
        conclusion="複数の論点を踏まえて次の判断へ進む。",
    )
    paired = dict(
        _base("paired_comparison"), left_label="現行", right_label="将来",
        rows=[
            {"criterion": f"評価軸{i + 1}", "left": "現行方式の特徴を記載する",
             "right": "将来方式の特徴を記載する"}
            for i in range(6)
        ],
        takeaway="評価軸ごとの差を踏まえて採用方針を判断する。",
    )
    left = [{"id": f"l{i}", "text": f"課題{i + 1}"} for i in range(6)]
    right = [{"id": f"r{i}", "text": f"施策{i + 1}"} for i in range(6)]
    links = [{"from": f"l{i}", "to": f"r{i}"} for i in range(6)]
    links.extend([
        {"from": "l0", "to": "r2"}, {"from": "l1", "to": "r4"},
        {"from": "l3", "to": "r5"}, {"from": "l5", "to": "r1"},
    ])
    mapping = dict(
        _base("mapping"), left_label="課題", right_label="施策",
        left_items=left, right_items=right, links=links,
    )
    lanes = [{"id": f"lane{i}", "label": f"担当{i + 1}"} for i in range(6)]
    stages = [{"id": f"stage{i}", "label": f"段階{i + 1}"} for i in range(6)]
    steps = [
        {"id": f"step{i}", "name": f"作業{i + 1}",
         "lane": f"lane{i % 6}", "stage": f"stage{i // 2}"}
        for i in range(12)
    ]
    edges = [{"from": f"step{i}", "to": f"step{i + 1}"} for i in range(11)]
    edges.append({"from": "step11", "to": "step0", "kind": "feedback"})
    swimlane = dict(
        _base("swimlane"), lanes=lanes, stages=stages, steps=steps, edges=edges,
    )
    participants = [
        {"id": f"participant{i}", "label": f"関係者{i + 1}"}
        for i in range(6)
    ]
    messages = [
        {"id": f"message{i}", "from": f"participant{i % 6}",
         "to": f"participant{(i + 2) % 6}", "label": f"確認{i + 1}",
         "kind": "return" if i % 4 == 3 else "request"}
        for i in range(12)
    ]
    sequence = dict(
        _base("sequence"), participants=participants, messages=messages,
        phases=[
            {"label": "準備", "from": "message0", "to": "message3"},
            {"label": "実行", "from": "message4", "to": "message7"},
            {"label": "確認", "from": "message8", "to": "message11"},
        ],
    )
    return [scope, summary, paired, mapping, swimlane, sequence]


def _assert_fit_stages():
    args = dict(
        count=5, row_h=0.68, min_row_h=0.46,
        gap=0.08, min_gap=0.02, font=13.5, min_font=10.0,
    )
    assert _fit_rows("test", 4.0, **args).stage == "standard"
    assert _fit_rows("test", 3.6, **args).stage == "gap"
    assert _fit_rows("test", 3.3, **args).stage == "element"
    assert _fit_rows("test", 2.4, **args).stage == "font"
    try:
        _fit_rows("test", 2.2, **args)
    except FitError:
        pass
    else:
        raise AssertionError("最小値でも収まらない入力を明示停止できていません")


def _assert_mapping_order():
    expected_crossings = {"疎": 0, "標準": 0, "上限": 3, "長文": 0}
    specs = [spec for spec in REVIEW_DECK["slides"] if spec["type"] == "mapping"]
    for spec in specs:
        left, right, crossings = _mapping_items_by_min_crossings(
            spec["left_items"], spec["right_items"], spec["links"])
        variant = spec["kicker"].split("/")[-1].strip()
        assert crossings == expected_crossings[variant], (variant, crossings)
        assert {item["id"] for item in left} == {
            item["id"] for item in spec["left_items"]}
        assert {item["id"] for item in right} == {
            item["id"] for item in spec["right_items"]}


def _assert_sequence_structure():
    spec = deepcopy(next(
        spec for spec in REVIEW_DECK["slides"]
        if spec["type"] == "sequence" and "標準" in spec["kicker"]
    ))
    slide = _render(_presentation(), spec)
    texts = {
        shape.text.strip()
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    }
    redundant = {f"{index + 1:02d}" for index in range(len(spec["messages"]))}
    assert texts.isdisjoint(redundant), texts & redundant
    assert {phase["label"] for phase in spec["phases"]} <= texts
    participant_shapes = [
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text.strip() in {item["label"] for item in spec["participants"]}
    ]
    center = sum(
        shape.left + shape.width / 2 for shape in participant_shapes
    ) / len(participant_shapes)
    assert abs(center - _presentation().slide_width / 2) <= Inches(0.05)


def _assert_swimlane_legend():
    specs = {
        spec["kicker"].split("/")[-1].strip(): deepcopy(spec)
        for spec in REVIEW_DECK["slides"]
        if spec["type"] == "swimlane"
    }
    legend_texts = {"凡例", "順方向", "差戻し"}
    dense_slide = _render(_presentation(), specs["上限"])
    dense_texts = {
        shape.text.strip()
        for shape in dense_slide.shapes
        if getattr(shape, "has_text_frame", False)
    }
    assert legend_texts <= dense_texts
    standard_slide = _render(_presentation(), specs["標準"])
    standard_texts = {
        shape.text.strip()
        for shape in standard_slide.shapes
        if getattr(shape, "has_text_frame", False)
    }
    assert legend_texts.isdisjoint(standard_texts)


def main():
    errors = validate(deepcopy(PATTERN_DECK), allow_sample_content=True)
    assert not errors, "\n".join(errors)
    review_errors = validate(deepcopy(REVIEW_DECK), allow_sample_content=True)
    assert not review_errors, "\n".join(review_errors)

    samples = [
        deepcopy(spec) for spec in PATTERN_DECK["slides"]
        if spec["type"] in RENDERERS
    ]
    assert {spec["type"] for spec in samples} == set(RENDERERS)
    prs = _presentation()
    review_specs = [deepcopy(spec) for spec in REVIEW_DECK["slides"]]
    assert len(review_specs) == len(RENDERERS) * 4
    for spec in samples + review_specs + _dense_specs():
        slide = _render(prs, spec)
        _assert_in_slide(slide)

    bad_mapping = deepcopy(next(spec for spec in samples if spec["type"] == "mapping"))
    bad_mapping["links"][0]["to"] = "undefined"
    errors = validate({"meta": {"title": "検証"}, "slides": [bad_mapping]},
                      allow_sample_content=True)
    assert any("未定義id" in error for error in errors)
    _assert_fit_stages()
    _assert_mapping_order()
    _assert_sequence_structure()
    _assert_swimlane_legend()
    print("candidate renderer tests: OK")


if __name__ == "__main__":
    main()
