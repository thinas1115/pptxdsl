"""候補rendererの疎・標準・上限入力をまとめた目視確認デッキを生成する。"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import generate
from candidate_renderers import (
    s_mapping,
    s_paired_comparison,
    s_scope,
    s_sequence,
    s_summary,
    s_swimlane,
)
from candidate_review_cases import REVIEW_DECK
from validate_content import validate


RENDER = {
    "scope_boundary": s_scope,
    "decision_summary": s_summary,
    "paired_comparison": s_paired_comparison,
    "relationship_map": s_mapping,
    "swimlane_flow": s_swimlane,
    "message_sequence": s_sequence,
}


def main(out_path):
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    errors = validate(REVIEW_DECK, allow_sample_content=True)
    if errors:
        raise SystemExit("NG: 候補rendererレビュー\n  - " + "\n  - ".join(errors))
    generate.DECK = REVIEW_DECK
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    blank = prs.slide_layouts[6]
    for page, spec in enumerate(REVIEW_DECK["slides"], 1):
        slide = prs.slides.add_slide(blank)
        generate.render_slide(RENDER[spec["type"]], slide, spec, page)
        generate.footer(slide, page)
    prs.save(out_path)
    print(f"saved: {out_path} ({len(REVIEW_DECK['slides'])} slides)")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("out_path", nargs="?", default="out/candidate_review.pptx")
    args = parser.parse_args()
    main(args.out_path)
