"""社内テンプレート用のパターン検証ギャラリーを生成する。"""
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import generate
from content_patterns import PATTERN_DECK
from org_layout import s_org
from diagrams2 import s_matrix, s_process, s_program_roadmap
from diagram_layout import render_diagram
from image_slide import s_image
from candidate_renderers import (
    s_mapping, s_paired_comparison, s_scope, s_sequence, s_summary, s_swimlane,
)
from training_renderers import (
    s_code_lab, s_concept, s_knowledge_check, s_network, s_protocol_state_flow,
    s_protocol_anatomy,
)
from validate_content import validate


def s_diagram(slide, spec, page):
    """content.json と同じインライン仕様で構成図を描画する。"""
    area = generate.header(
        slide, spec["kicker"], spec["title"], spec.get("lead"))
    render_diagram(slide, spec["diagram"], note=spec.get("note"),
                   content_area=area if area.shifted else None)


RENDER = dict(generate.RENDER,
              org=s_org,
              process=s_process, program_roadmap=s_program_roadmap,
              matrix=s_matrix,
              diagram=s_diagram, image=s_image,
              scope=s_scope, summary=s_summary,
              paired_comparison=s_paired_comparison, mapping=s_mapping,
              swimlane=s_swimlane, sequence=s_sequence,
              concept=s_concept, network=s_network,
              protocol_state_flow=s_protocol_state_flow,
              protocol_anatomy=s_protocol_anatomy,
              code_lab=s_code_lab, knowledge_check=s_knowledge_check)


def main(out_path, cover_footer_config=None):
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        generate.configure_cover_footer(cover_footer_config)
    except ValueError as e:
        raise SystemExit(f"NG: 表紙・フッター設定: {e}") from e
    errors = validate(PATTERN_DECK, allow_sample_content=True)
    if errors:
        raise SystemExit("NG: パターンギャラリー\n  - " + "\n  - ".join(errors))
    # footer() は generate.DECK の meta を参照するため、ギャラリー用に差し替える。
    generate.DECK = PATTERN_DECK
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    blank = prs.slide_layouts[6]
    for idx, spec in enumerate(PATTERN_DECK["slides"], 1):
        slide = prs.slides.add_slide(blank)
        generate.render_slide(RENDER[spec["type"]], slide, spec, idx)
        if spec["type"] != "title":
            generate.footer(slide, idx)
    prs.save(out_path)
    print(f"saved: {out_path} ({len(PATTERN_DECK['slides'])} slides)")


if __name__ == "__main__":
    default_out = Path(__file__).resolve().parent.parent / "out" / "pattern_gallery.pptx"
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("out_path", nargs="?", default=default_out)
    parser.add_argument("--cover-footer-config", metavar="PATH",
                        help="表紙・フッター設定JSON")
    args = parser.parse_args()
    main(args.out_path, args.cover_footer_config)
