"""技術研修で使う概念説明、ネットワーク、プロトコル、実習renderer。"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from asset_paths import resolve_icon_path
from diagrams import add_arrow, arrow_label
from generate import (
    ACCENT,
    BODY_W,
    CANVAS,
    CORAL,
    GRAY,
    LIGHT,
    MARGIN,
    NAVY,
    RULE,
    TEXT,
    WHITE,
    ZEBRA,
    add_rect,
    add_text,
    header,
)
from layout_fit import FitError, fit_text_or_raise, select_fit, stepped
from textfit import line_height_in, text_width_in, wrap_compact, wrap_text


_LANE_PALETTE = (
    (RGBColor(0xD9, 0xEB, 0xE8), ACCENT),
    (RGBColor(0xE4, 0xE9, 0xEE), NAVY),
    (RGBColor(0xF1, 0xE3, 0xDF), CORAL),
    (RGBColor(0xEA, 0xE8, 0xE2), GRAY),
)
_FIELD_FILLS = {
    "standard": (WHITE, NAVY),
    "muted": (ZEBRA, GRAY),
    "highlight": (LIGHT, ACCENT),
    "alert": (RGBColor(0xF1, 0xE3, 0xDF), CORAL),
}


def _fit_concept(area, definition, points, has_misconception):
    left_w = 5.45
    right_w = BODY_W - left_w - 0.72
    available = area.height - 0.34 - (0.72 if has_misconception else 0.08)

    def used(values):
        definition_h = len(wrap_text(
            definition, left_w - 0.10, values["definition_font"], "regular"
        )) * line_height_in(values["definition_font"], 1.18)
        left_h = 0.96 + definition_h
        right_h = 0.0
        for point in points:
            label_h = len(wrap_text(
                point["label"], right_w, values["label_font"], "bold"
            )) * line_height_in(values["label_font"], 1.08)
            text_h = len(wrap_text(
                point["text"], right_w, values["body_font"], "regular"
            )) * line_height_in(values["body_font"], 1.14)
            right_h += label_h + text_h + values["inside_gap"]
        right_h += max(0, len(points) - 1) * values["row_gap"]
        return max(left_h, right_h)

    def candidates():
        standard = {
            "definition_font": 18.5, "label_font": 13.0,
            "body_font": 15.0, "inside_gap": 0.12, "row_gap": 0.28,
        }
        yield "standard", standard, used(standard)
        for row_gap in stepped(0.22, 0.08, 0.035):
            values = dict(standard, row_gap=row_gap)
            yield "gap", values, used(values)
        for body_font in stepped(14.5, 12.5, 0.5):
            values = {
                "definition_font": max(17.0, body_font + 4.5),
                "label_font": max(11.5, body_font - 2.0),
                "body_font": body_font,
                "inside_gap": 0.08,
                "row_gap": 0.08,
            }
            yield "font", values, used(values)

    return select_fit(
        "concept", available, candidates(),
        guidance="要点を4件以内へ絞るか、定義と説明を短くしてください。",
    )


def s_concept(slide, spec, page):
    """技術用語を定義し、要点と誤解を一続きの視線で説明する。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    points = spec["points"]
    misconception = spec.get("misconception")
    fitted = _fit_concept(area, spec["definition"], points, bool(misconception))
    values = fitted.values

    top = area.top + 0.34
    left_w = 5.45
    divider_x = MARGIN + left_w + 0.34
    right_x = divider_x + 0.38
    right_w = MARGIN + BODY_W - right_x
    content_bottom = area.bottom - (0.78 if misconception else 0.10)

    term_x = MARGIN + 0.04
    term_w = left_w - 0.10
    definition_x = MARGIN + 0.04
    definition_w = left_w - 0.10
    if spec.get("icon"):
        icon_path = resolve_icon_path(spec["icon"])
        icon_size = 0.60
        slide.shapes.add_picture(
            str(icon_path), Inches(MARGIN + 0.04), Inches(top + 0.06),
            Inches(icon_size), Inches(icon_size),
        )
        term_x += 0.84
        term_w -= 0.84
    term_size, term_lines = fit_text_or_raise(
        "concept", "term", spec["term"], term_w, 0.72, 30.0,
        min_pt=23.5, weight="bold", spacing=1.0, role="natural",
    )
    add_text(
        slide, term_x, top, term_w, 0.72, "\n".join(term_lines), term_size,
        bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0,
    )

    definition_y = top + 0.96
    definition_h = content_bottom - definition_y
    definition_size, definition_lines = fit_text_or_raise(
        "concept", "definition", spec["definition"], definition_w,
        definition_h, values["definition_font"], min_pt=16.5,
        spacing=1.18, role="natural",
    )
    add_text(
        slide, definition_x, definition_y, definition_w, definition_h,
        "\n".join(definition_lines), definition_size, color=TEXT,
        spacing=1.18,
    )

    add_rect(slide, divider_x, top, 0.01, content_bottom - top, RULE)
    row_h = (content_bottom - top) / len(points)
    for index, point in enumerate(points):
        y = top + index * row_h
        label_size, label_lines = fit_text_or_raise(
            "concept", f"points[{index}].label", point["label"],
            right_w, 0.34, values["label_font"], min_pt=11.5,
            weight="bold", spacing=1.08, role="compact",
        )
        label_h = len(label_lines) * line_height_in(label_size, 1.08) + 0.02
        add_text(
            slide, right_x, y + 0.03, right_w, label_h,
            "\n".join(label_lines), label_size, bold=True, color=ACCENT,
            spacing=1.08,
        )
        body_y = y + label_h + values["inside_gap"]
        body_size, body_lines = fit_text_or_raise(
            "concept", f"points[{index}].text", point["text"],
            right_w, row_h - (body_y - y) - 0.08, values["body_font"],
            min_pt=12.5, spacing=1.14, role="body",
        )
        add_text(
            slide, right_x, body_y, right_w, row_h - (body_y - y) - 0.08,
            "\n".join(body_lines), body_size, color=TEXT, spacing=1.14,
        )
        if index < len(points) - 1:
            add_rect(slide, right_x, y + row_h - 0.02, right_w, 0.008, RULE)

    if misconception:
        y = area.bottom - 0.63
        add_rect(slide, MARGIN, y, BODY_W, 0.01, RULE)
        add_text(slide, MARGIN + 0.04, y + 0.13, 1.18, 0.24,
                 "誤解しやすい点", 10.5, bold=True, color=CORAL)
        size, lines = fit_text_or_raise(
            "concept", "misconception", misconception, BODY_W - 1.42,
            0.38, 13.0, min_pt=11.0, spacing=1.10, role="compact",
        )
        add_text(slide, MARGIN + 1.42, y + 0.10, BODY_W - 1.42, 0.38,
                 "\n".join(lines), size, color=TEXT, spacing=1.10)


def _set_connector_arrow(connector, *, both=False, dash=None):
    ln = connector.line._get_or_add_ln()
    for tag, enabled in (("a:tailEnd", True), ("a:headEnd", both)):
        if enabled:
            ln.append(ln.makeelement(
                qn(tag), {"type": "triangle", "w": "med", "len": "med"}))
    if dash:
        ln.insert(0, ln.makeelement(qn("a:prstDash"), {"val": dash}))


def _plain_connector(slide, x1, y1, x2, y2, color, *, width=1.6,
                     arrow=False, dash=None):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.shadow.inherit = False
    if arrow:
        _set_connector_arrow(conn, dash=dash)
    elif dash:
        ln = conn.line._get_or_add_ln()
        ln.insert(0, ln.makeelement(qn("a:prstDash"), {"val": dash}))
    return conn


def _route_colored(slide, points, color, *, width=1.6, dash=None):
    for start, end in zip(points[:-2], points[1:-1]):
        _plain_connector(slide, *start, *end, color, width=width, dash=dash)
    add_arrow(slide, *points[-2], *points[-1], color=color,
              width=width, dash=dash)


def _compact_lane_names(labels):
    parts = [label.split(maxsplit=1) for label in labels]
    if all(len(part) == 2 for part in parts):
        prefix = parts[0][0]
        if all(part[0] == prefix for part in parts[1:]):
            return f"{prefix} " + " / ".join(part[1] for part in parts)
    return " / ".join(labels)


def _fit_network(area, lane_count, max_cell_nodes):
    available = area.height - 0.58

    def used(values):
        return lane_count * values["lane_h"] + max(0, lane_count - 1) * values["gap"]

    def candidates():
        standard_lane = 1.62 if lane_count == 1 else 1.42 if lane_count == 2 else 1.32
        standard_icon = 0.48 if max_cell_nodes > 1 else 0.62 if lane_count == 1 else 0.56
        standard = {"lane_h": standard_lane,
                    "gap": 0.10, "icon": standard_icon,
                    "font": 9.5 if max_cell_nodes > 1 else 10.5}
        yield "standard", standard, used(standard)
        for gap in stepped(0.07, 0.02, 0.025):
            values = dict(standard, gap=gap)
            yield "gap", values, used(values)
        min_lane = 0.90 if lane_count >= 3 else 1.02
        for lane_h in stepped(standard["lane_h"] - 0.05, min_lane, 0.05):
            values = {"lane_h": lane_h, "gap": 0.02,
                      "icon": 0.44 if max_cell_nodes > 1 else 0.50,
                      "font": 9.0 if max_cell_nodes > 1 else 10.0}
            yield "element", values, used(values)
        values = {"lane_h": min_lane, "gap": 0.02, "icon": 0.44,
                  "font": 9.0}
        yield "font", values, used(values)

    return select_fit(
        "network", available, candidates(),
        guidance="論理セグメントまたは機器数を減らし、複数スライドへ分割してください。")


def _node_label(slide, x, y, width, text, font, *, sub=None):
    title_size, title_lines = fit_text_or_raise(
        "network", "node.label", text, width, 0.35, font,
        min_pt=8.5, weight="bold", spacing=1.05, role="compact")
    rendered = "\n".join(title_lines)
    title_h = len(title_lines) * line_height_in(title_size, 1.05) + 0.02
    box = add_text(
        slide, x - width / 2, y, width, title_h, rendered, title_size,
        bold=True, color=NAVY, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)
    box.fill.solid()
    box.fill.fore_color.rgb = CANVAS
    if sub:
        sub_size, sub_lines = fit_text_or_raise(
            "network", "node.sub", sub, width, 0.28, font - 1.5,
            min_pt=7.5, spacing=1.0, role="compact")
        sub_h = len(sub_lines) * line_height_in(sub_size, 1.0) + 0.01
        sub_box = add_text(
            slide, x - width / 2, y + title_h + 0.01, width, sub_h,
            "\n".join(sub_lines), sub_size, color=GRAY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
        sub_box.fill.solid()
        sub_box.fill.fore_color.rgb = CANVAS


def s_network(slide, spec, page):
    """物理機器と論理セグメントを同じ図で読めるネットワーク図。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    lanes = spec["lanes"]
    columns = spec["columns"]
    nodes = spec["nodes"]
    links = spec["links"]
    lane_index = {lane["id"]: index for index, lane in enumerate(lanes)}
    placement_groups = {}
    for node in nodes:
        ordered_lanes = tuple(sorted(node["lanes"], key=lane_index.get))
        placement_groups.setdefault((node["column"], ordered_lanes), []).append(node["id"])
    max_cell_nodes = max(len(group) for group in placement_groups.values())
    if max_cell_nodes > 2:
        raise FitError(
            "network.nodes: 同じ列・論理セグメントに3件以上の機器は配置できません。"
            "列を追加するか、複数スライドへ分割してください。")
    fitted = _fit_network(area, len(lanes), max_cell_nodes)
    lane_h = fitted.values["lane_h"]
    gap = fitted.values["gap"]
    icon_size = fitted.values["icon"]
    font = fitted.values["font"]

    top = area.top + 0.46
    used_h = len(lanes) * lane_h + max(0, len(lanes) - 1) * gap
    top += max(0.0, (area.height - 0.55 - used_h) * 0.12)
    lane_label_w = 1.46
    network_x0 = MARGIN + lane_label_w
    network_x1 = MARGIN + BODY_W - 0.16
    network_w = network_x1 - network_x0
    col_gap = network_w / len(columns)
    lane_by_id = {lane["id"]: lane for lane in lanes}
    lane_centers = {
        lane["id"]: top + index * (lane_h + gap) + lane_h / 2
        for index, lane in enumerate(lanes)
    }
    col_centers = {
        column["id"]: network_x0 + (index + 0.5) * col_gap
        for index, column in enumerate(columns)
    }
    col_index = {column["id"]: index for index, column in enumerate(columns)}

    for index, lane in enumerate(lanes):
        y = top + index * (lane_h + gap)
        fill, color = _LANE_PALETTE[index]
        background = add_rect(slide, MARGIN, y, BODY_W, lane_h, fill)
        background.name = f"layout-background:network-lane:{lane['id']}"
        add_text(slide, MARGIN + 0.16, y + 0.20, lane_label_w - 0.26, 0.27,
                 lane["label"], 11.5, bold=True, color=color,
                 anchor=MSO_ANCHOR.MIDDLE)
        if lane.get("sub"):
            add_text(slide, MARGIN + 0.16, y + 0.49, lane_label_w - 0.26, 0.22,
                     lane["sub"], 8.5, color=GRAY)

    for column in columns:
        if column.get("label"):
            add_text(slide, col_centers[column["id"]] - col_gap / 2,
                     top - 0.33, col_gap, 0.22, column["label"], 9,
                     bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    positions = {}
    node_by_id = {}
    placement_offsets = {}
    for group in placement_groups.values():
        if len(group) == 1:
            placement_offsets[group[0]] = 0.0
            continue
        spacing = min(0.72, col_gap * 0.42)
        placement_offsets[group[0]] = -spacing / 2
        placement_offsets[group[1]] = spacing / 2

    for node in nodes:
        lane_ids = node["lanes"]
        x = col_centers[node["column"]] + placement_offsets[node["id"]]
        y = sum(lane_centers[lane] for lane in lane_ids) / len(lane_ids)
        positions[node["id"]] = (x, y)
        node_by_id[node["id"]] = node

    def node_edge(node_id, toward_x, *, lane_id=None):
        x, y = positions[node_id]
        node = node_by_id[node_id]
        if lane_id and lane_id in node["lanes"]:
            y = lane_centers[lane_id]
        delta = icon_size / 2 + 0.07
        if len(node["lanes"]) > 1 and abs(y - positions[node_id][1]) > icon_size / 2:
            return x, y
        return (x + delta if toward_x > x else x - delta, y)

    for node in nodes:
        if len(node["lanes"]) <= 1:
            continue
        x, _y = positions[node["id"]]
        ordered = sorted(node["lanes"], key=lane_index.get)
        node_y = positions[node["id"]][1]
        rail_top = lane_centers[ordered[0]]
        rail_bottom = lane_centers[ordered[-1]]
        clearance = icon_size / 2 + 0.09
        if rail_top < node_y - clearance:
            _plain_connector(slide, x, rail_top, x, node_y - clearance,
                             RULE, width=2.0)
        if node_y + clearance < rail_bottom:
            _plain_connector(slide, x, node_y + clearance, x, rail_bottom,
                             RULE, width=2.0)
        for lane_id in ordered:
            port = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - 0.045), Inches(lane_centers[lane_id] - 0.045),
                Inches(0.09), Inches(0.09))
            port.fill.solid()
            port.fill.fore_color.rgb = NAVY
            port.line.fill.background()
            port.shadow.inherit = False

    # A broadcast is one frame copied to multiple ports. Draw those copies as
    # one shared branch instead of unrelated point-to-point arrows.
    broadcast_groups = {}
    for link_index, link in enumerate(links):
        if link.get("kind") != "broadcast":
            continue
        lane_ids = link.get("lanes", [])
        lane_id = lane_ids[0] if lane_ids else None
        broadcast_groups.setdefault((link["from"], lane_id), []).append(
            (link_index, link))

    grouped_broadcasts = set()
    for (source_id, lane_id), group in broadcast_groups.items():
        if len(group) < 2:
            continue
        target_ids = [link["to"] for _index, link in group]
        source_x, source_y = positions[source_id]
        target_xs = [positions[target_id][0] for target_id in target_ids]
        direction_x = sum(target_xs) / len(target_xs)
        start = node_edge(source_id, direction_x, lane_id=lane_id)
        ends = [node_edge(target_id, source_x, lane_id=lane_id)
                for target_id in target_ids]
        channel_y = source_y - lane_h * 0.30
        bus_end_x = max(end[0] for end in ends)
        _plain_connector(slide, start[0], start[1], start[0], channel_y,
                         CORAL, width=2.2)
        _plain_connector(slide, start[0], channel_y, bus_end_x, channel_y,
                         CORAL, width=2.2)
        for end in ends:
            add_arrow(slide, end[0], channel_y, end[0], end[1],
                      color=CORAL, width=2.2)
        labels = [link.get("label") for _index, link in group
                  if link.get("label")]
        if labels:
            arrow_label(
                slide, (start[0] + bus_end_x) / 2, channel_y - 0.16,
                labels[0],
                w=min(1.70, max(0.90, bus_end_x - start[0] - 0.10)),
                size=8.5)
        grouped_broadcasts.update(index for index, _link in group)

    # 配線を先に描き、アイコンとラベルを前面へ置く。
    legend_kinds = []
    for link_index, link in enumerate(links):
        source_x, source_y = positions[link["from"]]
        target_x, target_y = positions[link["to"]]
        kind = link.get("kind", "access")
        lane_ids = link.get("lanes", [])
        if kind not in legend_kinds:
            legend_kinds.append(kind)
        if link_index in grouped_broadcasts:
            continue
        if kind == "trunk":
            active = lane_ids or list(lane_by_id)
            for lane_id in active:
                color = _LANE_PALETTE[lane_index[lane_id]][1]
                start = node_edge(link["from"], target_x, lane_id=lane_id)
                end = node_edge(link["to"], source_x, lane_id=lane_id)
                _plain_connector(
                    slide, start[0], start[1], end[0], end[1],
                    color, width=2.0)
            lane_names = _compact_lane_names(
                [lane_by_id[lane_id]["label"] for lane_id in active])
            physical = link.get("label", "Trunk")
            trunk_label = f"{physical}\n{lane_names}"
            clear_w = max(0.80, abs(end[0] - start[0]) - 0.12)
            label_w = min(2.20, clear_w)
            label_size, label_lines = fit_text_or_raise(
                "network", "link.trunk.label", trunk_label,
                label_w - 0.12, 0.34, 8.5, min_pt=7.0,
                weight="bold", spacing=1.0, role="compact")
            label = add_text(
                slide, (source_x + target_x) / 2 - label_w / 2,
                min(source_y, target_y) - 0.42, label_w, 0.34,
                "\n".join(label_lines), label_size, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                spacing=1.0)
            label.fill.solid()
            label.fill.fore_color.rgb = CANVAS
            continue

        lane_id = lane_ids[0] if lane_ids else None
        color = (
            CORAL if kind == "broadcast" else
            _LANE_PALETTE[lane_index[lane_id]][1]
            if lane_id in lane_index else NAVY if kind == "routed" else GRAY
        )
        start = node_edge(link["from"], target_x, lane_id=lane_id)
        end = node_edge(link["to"], source_x, lane_id=lane_id)
        if kind == "blocked":
            boundary_y = (start[1] + end[1]) / 2
            stop_y = (boundary_y - 0.08 if end[1] > start[1]
                      else boundary_y + 0.08)
            middle_x = (start[0] + end[0]) / 2
            _plain_connector(slide, start[0], start[1], middle_x, start[1],
                             CORAL, width=1.8, dash="dash")
            _plain_connector(slide, middle_x, start[1], middle_x, stop_y,
                             CORAL, width=1.8, dash="dash")
            add_text(slide, middle_x - 0.19, boundary_y - 0.17, 0.38, 0.34,
                     "×", 18, bold=True, color=CORAL,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if link.get("label"):
                arrow_label(slide, middle_x + 0.78, boundary_y - 0.14,
                            link["label"], w=1.48, size=8.5)
            continue
        source_col = col_index[node_by_id[link["from"]]["column"]]
        target_col = col_index[node_by_id[link["to"]]["column"]]
        if (kind == "broadcast" and abs(source_col - target_col) > 1
                and abs(start[1] - end[1]) <= 0.02):
            channel_y = start[1] - lane_h * 0.30
            _route_colored(
                slide,
                [start, (start[0], channel_y), (end[0], channel_y), end],
                color, width=2.2,
            )
        elif abs(start[1] - end[1]) <= 0.02:
            add_arrow(slide, *start, *end, color=color,
                      width=2.2 if kind == "broadcast" else 1.6,
                      dash="dash" if kind == "control" else None)
        else:
            middle_x = (start[0] + end[0]) / 2
            _route_colored(slide, [start, (middle_x, start[1]),
                                   (middle_x, end[1]), end], color,
                           width=2.2 if kind == "broadcast" else 1.6,
                           dash="dash" if kind == "control" else None)
        if link.get("label"):
            arrow_label(
                slide, (start[0] + end[0]) / 2,
                (start[1] + end[1]) / 2 - 0.16,
                link["label"], w=min(1.4, col_gap * 0.88), size=8.5)

    for node in nodes:
        x, y = positions[node["id"]]
        halo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - icon_size / 2 - 0.06), Inches(y - icon_size / 2 - 0.06),
            Inches(icon_size + 0.12), Inches(icon_size + 0.12))
        halo.fill.solid()
        halo.fill.fore_color.rgb = CANVAS
        halo.line.fill.background()
        halo.shadow.inherit = False
        icon_path = resolve_icon_path(node["icon"])
        slide.shapes.add_picture(
            str(icon_path), Inches(x - icon_size / 2), Inches(y - icon_size / 2),
            Inches(icon_size), Inches(icon_size))
        group_size = len(placement_groups[(
            node["column"], tuple(sorted(node["lanes"], key=lane_index.get)))])
        label_width = min(1.68, col_gap * (0.44 if group_size > 1 else 0.92))
        _node_label(
            slide, x, y + icon_size / 2 + 0.08,
            label_width, node["label"], font,
            sub=node.get("sub"))

    legend_y = area.bottom - 0.25
    kind_labels = {"access": "Access", "trunk": "Trunk (複数VLAN)",
                   "routed": "L3転送", "control": "制御・確認",
                   "broadcast": "ブロードキャスト",
                   "blocked": "VLAN境界で停止"}
    x = MARGIN + 0.08
    for kind in legend_kinds:
        label = kind_labels[kind]
        width = text_width_in(label, 8.5) + 0.52
        color = (CORAL if kind in {"broadcast", "blocked"} else
                 ACCENT if kind in {"access", "trunk"} else NAVY)
        _plain_connector(slide, x, legend_y + 0.09, x + 0.28,
                         legend_y + 0.09, color,
                         width=2.2 if kind == "broadcast" else 2.0 if kind == "trunk" else 1.5,
                         arrow=kind in {"access", "routed", "broadcast"},
                         dash="dash" if kind in {"control", "blocked"} else None)
        add_text(slide, x + 0.36, legend_y, width - 0.36, 0.20,
                 label, 8.5, color=GRAY)
        x += width + 0.24


def _fit_protocol_state_flow(area, stages, flows):
    rail_w = 1.42
    available = area.height - 0.08

    def measured(values):
        stage_w = (
            BODY_W - rail_w - (len(stages) - 1) * values["stage_gap"]
        ) / len(stages)
        text_w = stage_w - 0.18
        header_lines = max(
            len(wrap_compact(stage["label"], text_w, values["stage_font"], "bold"))
            for stage in stages
        )
        header_h = max(
            values["icon_size"] + 0.08
            + header_lines * line_height_in(values["stage_font"], 1.0),
            0.62,
        )
        row_heights = []
        for flow in flows:
            flow_label_lines = wrap_compact(
                flow["label"], rail_w - 0.20, values["flow_font"], "bold"
            )
            flow_sub_lines = (
                wrap_compact(
                    flow["sub"], rail_w - 0.20,
                    values["flow_sub_font"], "regular"
                )
                if flow.get("sub") else []
            )
            rail_text_h = (
                0.20
                + len(flow_label_lines)
                * line_height_in(values["flow_font"], 1.0)
                + (0.04 + len(flow_sub_lines)
                   * line_height_in(values["flow_sub_font"], 1.0)
                   if flow_sub_lines else 0.0)
                + 0.14
            )
            state_heights = []
            for state in flow["states"]:
                title_lines = wrap_compact(
                    state["label"], text_w, values["state_font"], "bold"
                )
                detail_lines = wrap_compact(
                    state.get("detail", ""), text_w,
                    values["detail_font"], "regular"
                ) if state.get("detail") else []
                text_h = (
                    len(title_lines) * line_height_in(values["state_font"], 1.0)
                    + (0.05 + len(detail_lines)
                       * line_height_in(values["detail_font"], 1.0)
                       if detail_lines else 0.0)
                )
                state_heights.append(0.42 + text_h)
            row_heights.append(max(
                0.98, rail_text_h, max(state_heights) + 0.18
            ))
        total = (
            header_h + 0.18 + sum(row_heights)
            + max(0, len(flows) - 1) * values["row_gap"]
            + 0.06
        )
        result = dict(values)
        result.update({
            "stage_w": stage_w,
            "header_h": header_h,
            "row_heights": row_heights,
        })
        return result, total

    def candidates():
        standard = {
            "stage_gap": 0.18, "row_gap": 0.16,
            "icon_size": 0.34, "stage_font": 10.5,
            "flow_font": 12.5, "flow_sub_font": 9.0,
            "state_font": 11.0, "detail_font": 9.0,
        }
        values, used = measured(standard)
        yield "standard", values, used
        for row_gap in stepped(0.12, 0.06, 0.03):
            compact = dict(standard, row_gap=row_gap, stage_gap=0.12)
            values, used = measured(compact)
            yield "gap", values, used
        for icon_size in stepped(0.32, 0.26, 0.02):
            compact = dict(
                standard, row_gap=0.06, stage_gap=0.10,
                icon_size=icon_size,
            )
            values, used = measured(compact)
            yield "icon", values, used
        for state_font in stepped(10.5, 9.5, 0.5):
            compact = {
                "stage_gap": 0.10, "row_gap": 0.06,
                "icon_size": 0.26, "stage_font": 9.5,
                "flow_font": max(11.0, state_font + 1.0),
                "flow_sub_font": 8.0,
                "state_font": state_font,
                "detail_font": max(8.0, state_font - 2.0),
            }
            values, used = measured(compact)
            yield "font", values, used

    return select_fit(
        "protocol_state_flow", available, candidates(),
        guidance=(
            "段階を6件以内、トラックを3件以内に絞り、各状態の説明を短くするか、"
            "複数ページへ分割してください。"
        ),
    )


def _protocol_state_marker(
        slide, cx, y, width, appearance, color, encapsulation=None):
    marker_w = min(1.02, width * 0.72)
    marker_h = 0.20
    x = cx - marker_w / 2
    if appearance == "internal":
        marker = add_rect(slide, x, y + 0.08, marker_w, 0.04, color)
        marker.name = "protocol-state-flow:marker:internal"
        return marker_w
    fill = RGBColor(0xF9, 0xF8, 0xF4)
    line = CORAL if appearance == "alert" else color
    marker = add_rect(slide, x, y, marker_w, marker_h, fill, line=line)
    marker.name = f"protocol-state-flow:marker:{appearance}"
    if appearance == "encapsulated":
        tag_w = min(
            marker_w * 0.48,
            max(0.20, text_width_in(encapsulation, 6.5, "bold") + 0.10),
        )
        tag = add_rect(slide, x, y, tag_w, marker_h, color)
        tag.name = "protocol-state-flow:marker:encapsulated:tag"
        label_size, label_lines = fit_text_or_raise(
            "protocol_state_flow", "state.encapsulation", encapsulation,
            tag_w - 0.04, marker_h - 0.02, 6.5, min_pt=5.5,
            weight="bold", spacing=1.0, role="compact",
        )
        add_text(
            slide, x + 0.02, y + 0.01, tag_w - 0.04, marker_h - 0.02,
            "\n".join(label_lines), label_size,
            bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, spacing=1.0,
        )
    return marker_w


def s_protocol_state_flow(slide, spec, page):
    """端末・装置・伝送区間を通るフレームまたはパケットの状態変化を比較表示する。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    stages = spec["stages"]
    flows = spec["flows"]
    fitted = _fit_protocol_state_flow(area, stages, flows)
    values = fitted.values

    rail_w = 1.42
    stage_x0 = MARGIN + rail_w
    stage_w = values["stage_w"]
    stage_centers = {
        stage["id"]: (
            stage_x0 + index * (stage_w + values["stage_gap"]) + stage_w / 2
        )
        for index, stage in enumerate(stages)
    }
    top = area.top + 0.08
    rows_top = top + values["header_h"] + 0.18
    rows_bottom = (
        rows_top + sum(values["row_heights"])
        + max(0, len(flows) - 1) * values["row_gap"]
    )

    for stage in stages:
        cx = stage_centers[stage["id"]]
        if stage.get("role") == "link":
            background = add_rect(
                slide, cx - stage_w / 2 - 0.03, rows_top - 0.08,
                stage_w + 0.06, rows_bottom - rows_top + 0.16,
                RGBColor(0xEA, 0xF0, 0xEF),
            )
            background.name = (
                f"layout-background:protocol-state-flow:stage:{stage['id']}"
            )
        icon_path = resolve_icon_path(stage["icon"])
        icon_size = values["icon_size"]
        slide.shapes.add_picture(
            str(icon_path), Inches(cx - icon_size / 2), Inches(top),
            Inches(icon_size), Inches(icon_size),
        )
        label_size, label_lines = fit_text_or_raise(
            "protocol_state_flow", "stage.label", stage["label"],
            stage_w - 0.12, values["header_h"] - icon_size - 0.05,
            values["stage_font"], min_pt=9.0, weight="bold",
            spacing=1.0, role="compact",
        )
        add_text(
            slide, cx - stage_w / 2 + 0.06, top + icon_size + 0.06,
            stage_w - 0.12, values["header_h"] - icon_size - 0.04,
            "\n".join(label_lines), label_size, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, spacing=1.0,
        )

    row_y = rows_top
    for flow_index, flow in enumerate(flows):
        row_h = values["row_heights"][flow_index]
        fill, color = _LANE_PALETTE[flow_index % len(_LANE_PALETTE)]
        background = add_rect(slide, MARGIN, row_y, BODY_W, row_h, fill)
        background.name = (
            f"layout-background:protocol-state-flow:track:{flow_index}"
        )
        flow_label_size, flow_label_lines = fit_text_or_raise(
            "protocol_state_flow", "flow.label", flow["label"],
            rail_w - 0.20, row_h - 0.32,
            values["flow_font"], min_pt=10.5, weight="bold",
            spacing=1.0, role="compact",
        )
        flow_label_h = (
            len(flow_label_lines) * line_height_in(flow_label_size, 1.0) + 0.02
        )
        add_text(
            slide, MARGIN + 0.10, row_y + 0.20,
            rail_w - 0.20, flow_label_h,
            "\n".join(flow_label_lines), flow_label_size,
            bold=True, color=color, spacing=1.0,
        )
        if flow.get("sub"):
            flow_sub_y = row_y + 0.20 + flow_label_h + 0.04
            sub_size, sub_lines = fit_text_or_raise(
                "protocol_state_flow", "flow.sub", flow["sub"],
                rail_w - 0.20, row_y + row_h - flow_sub_y - 0.10,
                values["flow_sub_font"], min_pt=8.0,
                spacing=1.0, role="compact",
            )
            add_text(
                slide, MARGIN + 0.10, flow_sub_y,
                rail_w - 0.20, row_y + row_h - flow_sub_y - 0.10,
                "\n".join(sub_lines), sub_size, color=GRAY, spacing=1.0,
            )

        states = sorted(
            flow["states"],
            key=lambda state: next(
                index for index, stage in enumerate(stages)
                if stage["id"] == state["stage"]
            ),
        )
        marker_y = row_y + 0.18
        arrow_y = marker_y + 0.10
        marker_widths = {}
        for state in states:
            marker_widths[state["stage"]] = min(1.02, stage_w * 0.72)
        for left, right in zip(states, states[1:]):
            left_cx = stage_centers[left["stage"]]
            right_cx = stage_centers[right["stage"]]
            add_arrow(
                slide,
                left_cx + marker_widths[left["stage"]] / 2 + 0.04,
                arrow_y,
                right_cx - marker_widths[right["stage"]] / 2 - 0.04,
                arrow_y,
                color=color, width=1.6,
            )

        for state in states:
            cx = stage_centers[state["stage"]]
            _protocol_state_marker(
                slide, cx, marker_y, stage_w,
                state.get("appearance", "plain"), color,
                state.get("encapsulation"),
            )
            text_w = stage_w - 0.18
            title_size, title_lines = fit_text_or_raise(
                "protocol_state_flow", "state.label", state["label"],
                text_w, row_h - 0.48, values["state_font"],
                min_pt=9.5, weight="bold", spacing=1.0, role="compact",
            )
            title_y = marker_y + 0.30
            title_h = len(title_lines) * line_height_in(title_size, 1.0) + 0.03
            add_text(
                slide, cx - text_w / 2, title_y, text_w, title_h,
                "\n".join(title_lines), title_size, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER, spacing=1.0,
            )
            if state.get("detail"):
                detail_y = title_y + title_h + 0.03
                detail_h = row_y + row_h - detail_y - 0.10
                detail_size, detail_lines = fit_text_or_raise(
                    "protocol_state_flow", "state.detail", state["detail"],
                    text_w, detail_h, values["detail_font"],
                    min_pt=8.0, spacing=1.0, role="compact",
                )
                add_text(
                    slide, cx - text_w / 2, detail_y, text_w, detail_h,
                    "\n".join(detail_lines), detail_size, color=GRAY,
                    align=PP_ALIGN.CENTER, spacing=1.0,
                )
        row_y += row_h + values["row_gap"]

def _fit_protocol(area, frames, annotation_count):
    available = area.height
    annotation_rows = math.ceil(annotation_count / 4)

    def used(values):
        frame_block = (len(frames) * values["frame_h"]
                       + max(0, len(frames) - 1) * values["gap"])
        annotation_block = (
            0.28 + annotation_rows * values["annotation_row_h"]
            if annotation_rows else 0.0
        )
        return 0.30 + frame_block + annotation_block + 0.08

    def candidates():
        for gap in stepped(0.28, 0.12, 0.04):
            values = {"frame_h": 0.92, "gap": gap,
                      "annotation_row_h": 0.66, "annotation_text_h": 0.48,
                      "font": 10.5, "annotation_font": 10.5}
            yield "standard" if gap == 0.28 else "gap", values, used(values)
        for frame_h in stepped(0.86, 0.70, 0.04):
            values = {"frame_h": frame_h, "gap": 0.12,
                      "annotation_row_h": 0.58, "annotation_text_h": 0.42,
                      "font": 9.5, "annotation_font": 9.5}
            yield "element", values, used(values)
        values = {"frame_h": 0.68, "gap": 0.10,
                  "annotation_row_h": 0.54, "annotation_text_h": 0.40,
                  "font": 8.5, "annotation_font": 8.5}
        yield "font", values, used(values)

    return select_fit(
        "protocol_anatomy", available, candidates(),
        guidance="フレーム数・フィールド数・注釈数を減らして分割してください。")


def s_protocol_anatomy(slide, spec, page):
    """フレームやパケットのフィールド構造を比較する。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    frames = spec["frames"]
    annotations = [
        (frame_index, annotation)
        for frame_index, frame in enumerate(frames)
        for annotation in frame.get("annotations", [])
    ]
    fitted = _fit_protocol(area, frames, len(annotations))
    values = fitted.values
    frame_h, gap, font = values["frame_h"], values["gap"], values["font"]
    label_w = 1.25
    frame_totals = [
        sum(field["bits"] for field in frame["fields"])
        for frame in frames
    ]
    min_total = min(frame_totals)
    max_total = max(frame_totals)
    compare_lengths = len(frames) > 1 and min_total != max_total
    length_label_w = 0.78 if compare_lengths else 0.0
    bar_x = MARGIN + label_w
    bar_w = BODY_W - label_w - length_label_w
    top = area.top + 0.30

    minimum_by_bits = {}
    for frame in frames:
        for field in frame["fields"]:
            minimum = max(
                0.34,
                text_width_in(field["name"], font, "bold") + 0.24,
            )
            minimum_by_bits[field["bits"]] = max(
                minimum_by_bits.get(field["bits"], 0.0), minimum)

    def frame_width(frame, scale):
        return sum(
            max(minimum_by_bits[field["bits"]],
                scale * math.sqrt(field["bits"]))
            for field in frame["fields"]
        )

    minimum_required = max(frame_width(frame, 0.0) for frame in frames)
    if minimum_required > bar_w:
        raise FitError(
            "protocol_anatomy.fields: フィールド名の最小幅だけで描画領域を"
            "超えます。フィールド数を減らすか名称を短くしてください。")

    low, high = 0.0, bar_w
    for _ in range(40):
        mid = (low + high) / 2
        if max(frame_width(frame, mid) for frame in frames) <= bar_w:
            low = mid
        else:
            high = mid
    width_by_bits = {
        bits: max(minimum, low * math.sqrt(bits))
        for bits, minimum in minimum_by_bits.items()
    }

    field_centers = {}
    field_names = {}
    for frame_index, frame in enumerate(frames):
        y = top + frame_index * (frame_h + gap)
        add_text(slide, MARGIN + 0.02, y + 0.20, label_w - 0.18, 0.30,
                 frame["label"], 11, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        frame_bar_w = sum(
            width_by_bits[field["bits"]] for field in frame["fields"])
        x = bar_x
        for field_index, field in enumerate(frame["fields"]):
            width = width_by_bits[field["bits"]]
            fill, color = _FIELD_FILLS[field.get("role", "standard")]
            add_rect(slide, x, y, width, frame_h, fill, line=RULE)
            name_size, name_lines = fit_text_or_raise(
                "protocol_anatomy", f"frames[{frame_index}].fields[{field_index}].name",
                field["name"], width - 0.10, frame_h * 0.52, font,
                min_pt=7.5, weight="bold", spacing=1.0, role="compact")
            add_text(slide, x + 0.05, y + 0.12, width - 0.10, frame_h * 0.48,
                     "\n".join(name_lines), name_size, bold=True, color=color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
            add_text(slide, x + 0.04, y + frame_h - 0.25, width - 0.08, 0.16,
                     field.get("size_label", f"{field['bits']} bit"),
                     max(7.0, font - 2), color=GRAY,
                     align=PP_ALIGN.CENTER)
            field_centers[(frame_index, field["id"])] = (x + width / 2, y + frame_h)
            field_names[(frame_index, field["id"])] = field["name"]
            x += width

        if compare_lengths:
            delta = frame_totals[frame_index] - min_total
            if delta == 0:
                length_label = "基準"
            elif delta % 8 == 0:
                length_label = f"+{delta // 8} byte"
            else:
                length_label = f"+{delta} bit"
            add_text(slide, bar_x + frame_bar_w + 0.08, y + 0.33,
                     length_label_w - 0.10, 0.22, length_label,
                     max(7.5, font - 1.5), bold=delta > 0,
                     color=ACCENT if delta > 0 else GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)

    if annotations:
        annotation_y = top + len(frames) * frame_h + max(0, len(frames) - 1) * gap + 0.28
        columns = min(4, len(annotations))
        annotation_w = (BODY_W - 0.22 * (columns - 1)) / columns
        for index, (frame_index, annotation) in enumerate(annotations):
            col = index % columns
            row = index // columns
            x = MARGIN + col * (annotation_w + 0.22)
            y = annotation_y + row * values["annotation_row_h"]
            source_x, source_y = field_centers[(frame_index, annotation["field"])]
            target_x = x + 0.20
            target_y = y - 0.06
            _plain_connector(slide, source_x, source_y + 0.02,
                             target_x, target_y, RULE, width=0.8)
            annotation_label = field_names[(frame_index, annotation["field"])]
            annotation_label_w = min(
                0.62,
                max(0.42, text_width_in(annotation_label, 8.5, "bold") + 0.08),
            )
            add_text(slide, x, y, annotation_label_w, 0.22,
                     annotation_label, 8.5,
                     bold=True, color=ACCENT)
            fit_size, lines = fit_text_or_raise(
                "protocol_anatomy", "annotation", annotation["text"],
                annotation_w - annotation_label_w - 0.06,
                values["annotation_text_h"], values["annotation_font"],
                min_pt=8.0,
                spacing=1.08, role="compact")
            add_text(slide, x + annotation_label_w + 0.06, y,
                     annotation_w - annotation_label_w - 0.06,
                     values["annotation_text_h"],
                     "\n".join(lines), fit_size, color=TEXT, spacing=1.08)

def _fit_code_text(code, width, height):
    for size in stepped(13.0, 9.5, 0.5):
        widest = max(text_width_in(line or " ", size) for line in code.splitlines())
        used_h = max(1, len(code.splitlines())) * line_height_in(size, 1.18)
        if widest <= width and used_h <= height:
            return size
    raise FitError(
        "code_lab.code: 最小フォントでも折り返さずに収まりません。"
        "コマンドを短くするか、実習を複数スライドへ分割してください。")


def _apply_code_font(textbox):
    for paragraph in textbox.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Consolas"
            r_pr = run._r.get_or_add_rPr()
            for tag in ("a:ea", "a:cs"):
                element = r_pr.find(qn(tag))
                if element is None:
                    element = r_pr.makeelement(qn(tag), {})
                    r_pr.append(element)
                element.set("typeface", "Consolas")


def s_code_lab(slide, spec, page):
    """設定例と確認結果を、折り返さないコードと要点で示す。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    sections = spec["sections"]
    checks = spec["checks"]
    top = area.top + 0.30
    bottom = area.bottom - 0.15
    panel_h = bottom - top
    left_w = 7.55
    right_x = MARGIN + left_w + 0.42
    right_w = BODY_W - left_w - 0.42
    add_rect(slide, MARGIN, top, left_w, panel_h, NAVY)
    section_gap = 0.20
    section_h = (panel_h - section_gap * (len(sections) - 1)) / len(sections)
    for index, section in enumerate(sections):
        y = top + index * (section_h + section_gap)
        add_text(slide, MARGIN + 0.28, y + 0.20, left_w - 0.56, 0.22,
                 section["label"], 10, bold=True, color=RGBColor(0x8F, 0xD4, 0xCC))
        code_y = y + 0.55
        code_h = section_h - 0.72
        size = _fit_code_text(section["code"], left_w - 0.56, code_h)
        code_box = add_text(
            slide, MARGIN + 0.28, code_y, left_w - 0.56, code_h,
            section["code"], size, color=WHITE, spacing=1.18, wrap=False)
        _apply_code_font(code_box)

    add_text(slide, right_x, top + 0.02, right_w, 0.32,
             spec.get("check_label", "確認ポイント"), 13.5,
             bold=True, color=NAVY)
    rule_y = top + 0.47
    add_rect(slide, right_x, rule_y, right_w, 0.01, RULE)
    available = panel_h - 0.66
    row_h = available / len(checks)
    if row_h < 0.58:
        raise FitError("code_lab.checks: 確認項目が多すぎます。5件以内へ絞ってください。")
    for index, check in enumerate(checks):
        y = top + 0.62 + index * row_h
        add_text(slide, right_x, y + 0.02, 0.34, 0.24,
                 f"{index + 1:02d}", 9.5, bold=True, color=ACCENT)
        size, lines = fit_text_or_raise(
            "code_lab", f"checks[{index}]", check,
            right_w - 0.44, row_h - 0.10, 12.5, min_pt=9.5,
            spacing=1.13, role="body")
        add_text(slide, right_x + 0.44, y, right_w - 0.44, row_h - 0.08,
                 "\n".join(lines), size, color=TEXT, spacing=1.13)
        if index < len(checks) - 1:
            add_rect(slide, right_x + 0.44, y + row_h - 0.04,
                     right_w - 0.44, 0.008, RULE)

def _layout_options(question, size, mode, width):
    option_size = size - 1.5
    rows = []
    current = []
    current_width = 0.0
    for option_index, option in enumerate(question["options"]):
        correct = option_index == question["answer"]
        weight = "bold" if mode == "answers" and correct else "regular"
        option_text = f"{chr(65 + option_index)}. {option}"
        lines = wrap_compact(option_text, width, option_size, weight)
        natural_width = max(
            text_width_in(line or " ", option_size, weight) for line in lines
        ) + 0.28
        item_width = min(width, natural_width)
        item_height = max(
            0.28, len(lines) * line_height_in(option_size, 1.08) + 0.02
        )
        if current and current_width + 0.18 + item_width > width:
            rows.append(current)
            current = []
            current_width = 0.0
        current.append({
            "index": option_index,
            "text": "\n".join(lines),
            "width": item_width,
            "height": item_height,
        })
        current_width += (0.18 if len(current) > 1 else 0.0) + item_width
    if current:
        rows.append(current)
    height = sum(max(item["height"] for item in row) for row in rows)
    height += max(0, len(rows) - 1) * 0.05
    return rows, height


def _fit_questions(area, questions, mode):
    available = area.height - 0.36

    def measure(size, gap):
        heights = []
        for question in questions:
            question_h = len(wrap_text(question["question"], 10.25, size, "bold")) \
                * line_height_in(size, 1.12)
            _option_rows, option_h = _layout_options(
                question, size, mode, 10.85)
            explanation_h = 0.0
            if mode == "answers":
                explanation_h = len(wrap_text(
                    question["explanation"], 9.55, size - 2.0)) \
                    * line_height_in(size - 2.0, 1.10) + 0.10
            heights.append(max(0.76, question_h + option_h + explanation_h + 0.20))
        return heights, sum(heights) + max(0, len(heights) - 1) * gap

    def candidates():
        for gap in stepped(0.24, 0.10, 0.035):
            _heights, used = measure(15.0, gap)
            yield "standard" if gap == 0.24 else "gap", {"size": 15.0, "gap": gap}, used
        for size in stepped(14.5, 11.5, 0.5):
            _heights, used = measure(size, 0.10)
            yield "font", {"size": size, "gap": 0.10}, used

    fitted = select_fit(
        "knowledge_check", available, candidates(),
        guidance="設問を減らすか、選択肢と解説を短くしてください。")
    heights, _used = measure(fitted.values["size"], fitted.values["gap"])
    return fitted, heights


def s_knowledge_check(slide, spec, page):
    """設問用と解答・解説用を同じ構造から描き分ける。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    questions = spec["questions"]
    mode = spec["mode"]
    fitted, heights = _fit_questions(area, questions, mode)
    size, gap = fitted.values["size"], fitted.values["gap"]
    top = area.top + 0.22
    used = sum(heights) + gap * max(0, len(heights) - 1)
    top += max(0, (area.height - 0.30 - used) * 0.12)

    for index, (question, row_h) in enumerate(zip(questions, heights)):
        marker_color = ACCENT if mode == "answers" else NAVY
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(MARGIN + 0.02), Inches(top + 0.02),
            Inches(0.42), Inches(0.42))
        marker.fill.solid()
        marker.fill.fore_color.rgb = marker_color
        marker.line.fill.background()
        marker.shadow.inherit = False
        add_text(slide, MARGIN + 0.02, top + 0.02, 0.42, 0.42,
                 f"{index + 1:02d}", 10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = MARGIN + 0.68
        question_size, question_lines = fit_text_or_raise(
            "knowledge_check", f"questions[{index}].question",
            question["question"], 10.85, 0.56, size,
            min_pt=11.5, weight="bold", spacing=1.12, role="natural")
        q_h = len(question_lines) * line_height_in(question_size, 1.12)
        add_text(slide, x, top, 10.85, q_h + 0.04,
                 "\n".join(question_lines), question_size,
                 bold=True, color=NAVY, spacing=1.12)
        option_y = top + q_h + 0.10
        option_rows, option_h = _layout_options(
            question, size, mode, 10.85)
        for row in option_rows:
            option_x = x
            option_row_h = max(item["height"] for item in row)
            for item in row:
                correct = item["index"] == question["answer"]
                option_color = ACCENT if mode == "answers" and correct else GRAY
                add_text(
                    slide, option_x, option_y, item["width"], item["height"],
                    item["text"], size - 1.5,
                    bold=mode == "answers" and correct,
                    color=option_color, spacing=1.08,
                )
                option_x += item["width"] + 0.18
            option_y += option_row_h + 0.05
        if mode == "answers":
            explanation_y = top + q_h + 0.10 + option_h + 0.10
            add_text(slide, x, explanation_y, 0.58, 0.22,
                     "解説", 9.5, bold=True, color=ACCENT)
            exp_size, exp_lines = fit_text_or_raise(
                "knowledge_check", f"questions[{index}].explanation",
                question["explanation"], 9.90, row_h - (explanation_y - top),
                size - 2.0, min_pt=9.5, spacing=1.10, role="body")
            add_text(slide, x + 0.68, explanation_y, 9.90,
                     row_h - (explanation_y - top), "\n".join(exp_lines),
                     exp_size, color=TEXT, spacing=1.10)
        if index < len(questions) - 1:
            add_rect(slide, x, top + row_h + gap / 2,
                     BODY_W - 0.72, 0.008, RULE)
        top += row_h + gap
