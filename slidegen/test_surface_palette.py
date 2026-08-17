"""ギャラリーの通常面に真っ白な塗りが混入しないことを検証する。"""

from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from generate_patterns import main as generate_gallery


def _rgb_hex(fill):
    """単色塗りのRGBを返す。塗りなし・テーマ色・画像は対象外。"""
    try:
        rgb = fill.fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return str(rgb) if rgb is not None else None


def _is_near_white(rgb_hex):
    if not rgb_hex or len(rgb_hex) != 6:
        return False
    channels = [int(rgb_hex[index:index + 2], 16) for index in (0, 2, 4)]
    return min(channels) >= 250


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def assert_no_near_white_surfaces(path):
    prs = Presentation(path)
    findings = []
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows, 1):
                    for col_index, cell in enumerate(row.cells, 1):
                        rgb = _rgb_hex(cell.fill)
                        if _is_near_white(rgb):
                            findings.append(
                                f"{slide_index}枚目 表セル({row_index}, {col_index}): {rgb}")
                continue
            fill = getattr(shape, "fill", None)
            if fill is None:
                continue
            rgb = _rgb_hex(fill)
            if _is_near_white(rgb):
                findings.append(f"{slide_index}枚目 {shape.name}: {rgb}")
    assert not findings, (
        "通常面に真っ白な塗りを使用しないでください。"
        "濃色面上の白文字はこの検査の対象外です。\n  - "
        + "\n  - ".join(findings)
    )


def main():
    with TemporaryDirectory() as tmp:
        path = Path(tmp) / "pattern_gallery.pptx"
        generate_gallery(path)
        assert_no_near_white_surfaces(path)
    print("surface palette tests: OK")


if __name__ == "__main__":
    main()
