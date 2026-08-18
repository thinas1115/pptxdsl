"""大判画像を本文領域の主役として配置するrenderer。"""
from PIL import Image
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches

from asset_paths import resolve_image_path
from generate import (
    ACCENT, BODY_W, GRAY, LIGHT, ZEBRA, add_rect, add_text, header,
)
from layout_fit import select_fit, stepped

FRAME_X = 0.72
FRAME_W = BODY_W - 0.34
MIN_IMAGE_H = 2.40
PICTURE_NAME = "ContentImage"

COMPARE_GAP = 0.42
COMPARE_LEFT = 0.76
COMPARE_PANEL_W = (11.82 - COMPARE_GAP) / 2
MIN_COMPARE_IMAGE_H = 2.00
COMPARE_PICTURE_NAME = "ComparisonImage"


def fit_image_layout(available):
    """標準余白、余白圧縮、画像縮小、停止の順で収容する。"""
    candidates = []

    def add(stage, top_gap, bottom_gap, min_image_h):
        values = {
            "top_gap": top_gap,
            "bottom_gap": bottom_gap,
            "min_image_h": min_image_h,
        }
        candidates.append(
            (stage, values, top_gap + min_image_h + bottom_gap))

    add("standard", 0.06, 0.04, 3.60)
    add("gap", 0.02, 0.02, 3.20)
    for min_image_h in stepped(3.10, MIN_IMAGE_H, 0.10):
        add("element", 0.01, 0.01, min_image_h)
    return select_fit(
        "image", available, candidates,
        guidance=("leadを短くするか、画像だけのスライドへ分割してください。"),
    )


def _apply_offset_shadow(picture):
    """PowerPointの「オフセット: 右下」に近い外側影を設定する。"""
    sp_pr = picture._element.spPr
    existing = sp_pr.find(qn("a:effectLst"))
    if existing is not None:
        sp_pr.remove(existing)
    sp_pr.append(parse_xml(
        f'<a:effectLst {nsdecls("a")}>'
        '<a:outerShdw blurRad="76200" dist="50800" dir="2700000" '
        'algn="ctr" rotWithShape="0">'
        '<a:srgbClr val="000000"><a:alpha val="28000"/>'
        '</a:srgbClr></a:outerShdw></a:effectLst>'
    ))


def _add_picture(slide, path, x, y, w, h, fit, alt=None, shadow=False):
    with Image.open(path) as image:
        image_w, image_h = image.size
    image_ratio = image_w / image_h
    frame_ratio = w / h

    if fit == "cover":
        picture = slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), Inches(w), Inches(h))
        if image_ratio > frame_ratio:
            crop = (1 - frame_ratio / image_ratio) / 2
            picture.crop_left = picture.crop_right = crop
        elif image_ratio < frame_ratio:
            crop = (1 - image_ratio / frame_ratio) / 2
            picture.crop_top = picture.crop_bottom = crop
    else:
        if image_ratio >= frame_ratio:
            picture_w, picture_h = w, w / image_ratio
        else:
            picture_w, picture_h = h * image_ratio, h
        picture = slide.shapes.add_picture(
            str(path), Inches(x + (w - picture_w) / 2),
            Inches(y + (h - picture_h) / 2),
            Inches(picture_w), Inches(picture_h))

    picture.name = PICTURE_NAME
    if alt:
        picture._element.nvPicPr.cNvPr.set("descr", alt)
    if shadow:
        _apply_offset_shadow(picture)
    return picture


def s_image(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    path = resolve_image_path(spec["image"])
    if not path.is_file():
        raise FileNotFoundError(f"画像 {spec['image']!r} がassets内にありません")

    fitted = fit_image_layout(area.height)
    values = fitted.values
    top = area.top + values["top_gap"]
    image_h = area.height - values["top_gap"] - values["bottom_gap"]
    _add_picture(
        slide, path, FRAME_X, top, FRAME_W, image_h,
        spec.get("fit", "contain"), spec.get("alt"),
        spec.get("shadow", False))


def fit_image_compare_layout(available):
    """左右パネル共通の帯余白・画像最小高さを標準→圧縮→縮小の順で評価する。"""
    candidates = []

    def add(stage, top_pad, marker_h, mid_gap, bottom_pad, min_image_h):
        values = {
            "top_pad": top_pad, "marker_h": marker_h, "mid_gap": mid_gap,
            "bottom_pad": bottom_pad, "min_image_h": min_image_h,
        }
        used = top_pad + marker_h + mid_gap + min_image_h + bottom_pad
        candidates.append((stage, values, used))

    add("standard", 0.24, 0.30, 0.16, 0.20, 3.00)
    add("gap", 0.14, 0.30, 0.10, 0.12, 2.70)
    for min_image_h in stepped(2.60, MIN_COMPARE_IMAGE_H, 0.10):
        add("element", 0.10, 0.30, 0.08, 0.10, min_image_h)
    return select_fit(
        "image_compare", available, candidates,
        guidance=("leadを短くするか、片方の画像だけのimageスライドへ分割してください。"),
    )


def s_image_compare(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    panels = [spec["left"], spec["right"]]
    paths = []
    for side, panel in zip(("left", "right"), panels):
        path = resolve_image_path(panel["image"])
        if not path.is_file():
            raise FileNotFoundError(
                f"画像 {panel['image']!r} がassets内にありません")
        paths.append(path)

    fitted = fit_image_compare_layout(area.height)
    v = fitted.values
    panel_h = v["top_pad"] + v["marker_h"] + v["mid_gap"] \
        + v["min_image_h"] + v["bottom_pad"]
    top = area.top
    fit = spec.get("fit", "contain")
    shadow = spec.get("shadow", False)

    for i, (panel, path) in enumerate(zip(panels, paths)):
        x = COMPARE_LEFT + i * (COMPARE_PANEL_W + COMPARE_GAP)
        fill = ZEBRA if i == 0 else LIGHT
        marker = GRAY if i == 0 else ACCENT
        add_rect(slide, x, top, COMPARE_PANEL_W, panel_h, fill)
        add_text(slide, x + 0.30, top + v["top_pad"],
                 COMPARE_PANEL_W - 0.60, v["marker_h"],
                 panel.get("label", "BEFORE" if i == 0 else "AFTER"), 9.5,
                 bold=True, color=marker)
        image_top = top + v["top_pad"] + v["marker_h"] + v["mid_gap"]
        image_h = panel_h - v["top_pad"] - v["marker_h"] - v["mid_gap"] \
            - v["bottom_pad"]
        picture = _add_picture(
            slide, path, x + 0.20, image_top, COMPARE_PANEL_W - 0.40,
            image_h, fit, panel.get("alt"), shadow)
        picture.name = COMPARE_PICTURE_NAME
