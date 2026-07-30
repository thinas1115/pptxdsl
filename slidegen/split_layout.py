"""左右50:50の本文領域へ、対応済み子レイアウタを組み合わせる。"""

from PIL import Image
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

import generate
from asset_paths import resolve_image_path
from diagram_layout import render_diagram
from image_slide import _add_picture
from layout_fit import FitError, fit_text_or_raise, select_fit, stepped
from textfit import line_height_in, wrap_text


SUPPORTED_CHILD_TYPES = {
    "bullets", "cards", "table", "chart", "image", "diagram",
}
SPLIT_GAP = 0.72
SPLIT_LEFT = 0.72
SPLIT_RIGHT = 12.61
CHILD_HEADING_H = 0.68


def _regions(area, reserve_note):
    bottom = area.bottom - (0.30 if reserve_note else 0.0)
    width = (SPLIT_RIGHT - SPLIT_LEFT - SPLIT_GAP) / 2
    left = generate.ContentArea(
        area.top, bottom, area.shifted, SPLIT_LEFT, SPLIT_LEFT + width)
    right = generate.ContentArea(
        area.top, bottom, area.shifted,
        SPLIT_LEFT + width + SPLIT_GAP, SPLIT_RIGHT)
    return left, right


def _content_region(slide, side, child):
    region = child["_region"]
    inset = 0.04 if side == "left" else 0.28
    left = region.left + inset
    right = region.right - inset
    width = right - left
    number = "01" if side == "left" else "02"
    generate.add_text(
        slide, left, region.top + 0.04, 0.34, 0.24, number, 9.5,
        bold=True, color=generate.ACCENT, align=PP_ALIGN.LEFT)
    heading_size, heading_lines = fit_text_or_raise(
        "split", f"{side}.heading", child["heading"],
        width - 0.44, 0.38, 16.5,
        min_pt=13, weight="bold", spacing=1.08)
    generate.add_text(
        slide, left + 0.44, region.top + 0.01,
        width - 0.44, 0.38, "\n".join(heading_lines),
        heading_size, bold=True, color=generate.NAVY, spacing=1.08)
    return generate.ContentArea(
        region.top + CHILD_HEADING_H,
        region.bottom - 0.12,
        region.shifted,
        left,
        right,
    )


def _render_bullets(slide, side, child, area):
    bullets = child["bullets"]
    text_w = area.width - 0.52

    def measure(size, gap):
        heights = [
            len(wrap_text(item[0], text_w, size))
            * line_height_in(size, 1.18)
            for item in bullets
        ]
        return heights, sum(heights) + gap * max(0, len(heights) - 1)

    def candidates():
        for gap in stepped(0.30, 0.16, 0.02):
            _heights, used = measure(15, gap)
            yield ("standard" if gap == 0.30 else "gap",
                   {"size": 15, "gap": gap}, used)
        for size in stepped(14.5, 11.5, 0.5):
            _heights, used = measure(size, 0.16)
            yield "font", {"size": size, "gap": 0.16}, used

    fitted = select_fit(
        f"split.{side}.bullets", area.height - 0.10, candidates(),
        guidance="箇条書きを減らすか各項目を短くしてください。")
    size, gap = fitted.values["size"], fitted.values["gap"]
    heights, used = measure(size, gap)
    y = area.top + 0.08
    for index, (item, height) in enumerate(zip(bullets, heights), 1):
        generate.add_text(
            slide, area.left, y - 0.02, 0.34, 0.26, f"{index:02d}", 10,
            bold=True, color=generate.ACCENT, align=PP_ALIGN.RIGHT)
        generate.add_text(
            slide, area.left + 0.46, y, text_w, height + 0.05,
            item[0], size, color=generate.TEXT, spacing=1.18)
        y += height + gap


def _render_cards(slide, side, child, area):
    cards = [generate._normalize_card(card) for card in child["cards"]]
    text_w = area.width - 0.62

    def measure(size, gap):
        row_heights = []
        for card in cards:
            heading_lines = wrap_text(card["heading"], text_w, size, "bold")
            body_lines = wrap_text(card["body"], text_w, size - 1)
            row_heights.append(
                len(heading_lines) * line_height_in(size, 1.08)
                + len(body_lines) * line_height_in(size - 1, 1.15) + 0.18)
        return row_heights, sum(row_heights) + gap * max(0, len(cards) - 1)

    def candidates():
        for gap in stepped(0.22, 0.12, 0.02):
            _rows, used = measure(14, gap)
            yield ("standard" if gap == 0.22 else "gap",
                   {"size": 14, "gap": gap}, used)
        for size in stepped(13.5, 11, 0.5):
            _rows, used = measure(size, 0.12)
            yield "font", {"size": size, "gap": 0.12}, used

    fitted = select_fit(
        f"split.{side}.cards", area.height - 0.10, candidates(),
        guidance="カード本文を短くするかカード数を減らしてください。")
    size, gap = fitted.values["size"], fitted.values["gap"]
    row_heights, used = measure(size, gap)
    y = area.top + 0.06
    for index, (card, row_h) in enumerate(zip(cards, row_heights), 1):
        color = generate.CORAL if card["emphasis"] else generate.ACCENT
        generate.add_text(
            slide, area.left, y + 0.01, 0.36, 0.28, f"{index:02d}", 10,
            bold=True, color=color, align=PP_ALIGN.RIGHT)
        heading_size, heading_lines = fit_text_or_raise(
            f"split.{side}.cards", f"cards[{index - 1}].heading",
            card["heading"], text_w, row_h * 0.45, size,
            min_pt=11, weight="bold", spacing=1.08)
        heading_h = len(heading_lines) * line_height_in(heading_size, 1.08)
        generate.add_text(
            slide, area.left + 0.50, y, text_w, heading_h + 0.03,
            "\n".join(heading_lines), heading_size, bold=True,
            color=generate.NAVY, spacing=1.08)
        body_size, body_lines = fit_text_or_raise(
            f"split.{side}.cards", f"cards[{index - 1}].body",
            card["body"], text_w, row_h - heading_h - 0.10, size - 1,
            min_pt=10, spacing=1.15)
        generate.add_text(
            slide, area.left + 0.50, y + heading_h + 0.10,
            text_w, row_h - heading_h - 0.10, "\n".join(body_lines),
            body_size, color=generate.TEXT, spacing=1.15)
        if index < len(cards):
            generate.add_rect(
                slide, area.left + 0.50, y + row_h + gap / 2,
                text_w, 0.008, generate.RULE)
        y += row_h + gap


def _render_image(slide, side, child, area):
    path = resolve_image_path(child["image"])
    if not path.is_file():
        raise FileNotFoundError(f"画像 {child['image']!r} がassets内にありません")
    if area.height < 2.40:
        raise FitError(
            f"split.{side}.image: 最小画像高2.40inを確保できません。"
            "leadを短くするか画像を単独スライドへ移してください。")
    fit = child.get("fit", "contain")
    picture_h = area.height - 0.10
    if fit == "contain":
        with Image.open(path) as image:
            ratio = image.width / image.height
        picture_h = min(picture_h, area.width / ratio)
    _add_picture(
        slide, path, area.left, area.top + 0.05, area.width,
        picture_h, fit, child.get("alt"), child.get("shadow", False))


def _render_table(slide, side, child, area):
    columns, rows = child["columns"], child["rows"]
    col_w = area.width / len(columns)
    header_h = 0.54

    def measure(size, pad):
        row_hs = []
        for row in rows:
            needed = max(
                len(wrap_text(value, col_w - pad * 2, size))
                * line_height_in(size, 1.12)
                for value in row)
            row_hs.append(max(0.46, needed + pad * 2))
        return row_hs, header_h + sum(row_hs)

    def candidates():
        for pad in stepped(0.10, 0.06, 0.02):
            _rows, used = measure(10.5, pad)
            yield ("standard" if pad == 0.10 else "padding",
                   {"size": 10.5, "pad": pad}, used)
        for size in stepped(10, 8.5, 0.5):
            _rows, used = measure(size, 0.06)
            yield "font", {"size": size, "pad": 0.06}, used

    fitted = select_fit(
        f"split.{side}.table", area.height - 0.08, candidates(),
        guidance="表の行・列を減らすかセル内の文言を短くしてください。")
    size, pad = fitted.values["size"], fitted.values["pad"]
    row_hs, table_h = measure(size, pad)
    top = area.top + 0.04
    graphic = slide.shapes.add_table(
        len(rows) + 1, len(columns),
        Inches(area.left), Inches(top), Inches(area.width), Inches(table_h))
    table = graphic.table
    for column in table.columns:
        column.width = Emu(int(Inches(col_w)))
    table.rows[0].height = Emu(int(Inches(header_h)))
    for index, row_h in enumerate(row_hs, 1):
        table.rows[index].height = Emu(int(Inches(row_h)))
    for index, heading in enumerate(columns):
        generate._cell(
            table.cell(0, index), heading, size, bold=True,
            color=generate.WHITE, fill=generate.NAVY, center=index > 0)
    for row_index, row in enumerate(rows, 1):
        fill = generate.WHITE if row_index % 2 == 0 else generate.ZEBRA
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            generate._cell(
                cell, value, size, bold=col_index == 0,
                color=generate.NAVY if col_index == 0 else generate.TEXT,
                fill=fill, center=False)
            cell.margin_left = cell.margin_right = Inches(pad)


def _render_chart(slide, side, child, area):
    chart_spec = child["chart"]
    if area.height < 2.55:
        raise FitError(
            f"split.{side}.chart: グラフ高が不足しています。"
            "leadを短くするかグラフを単独スライドへ移してください。")
    chart_types = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
        "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    }
    kind = chart_spec.get("kind", "bar")
    data = CategoryChartData()
    data.categories = chart_spec["categories"]
    for name, values in chart_spec["series"]:
        data.add_series(name, values)
    graphic = slide.shapes.add_chart(
        chart_types[kind], Inches(area.left), Inches(area.top + 0.02),
        Inches(area.width), Inches(area.height - 0.08), data)
    chart = graphic.chart
    chart.has_title = False
    chart.has_legend = chart_spec.get(
        "show_legend", len(chart_spec["series"]) > 1)
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.font.name = generate.FONT
    chart.font.size = Pt(8.5)
    plot = chart.plots[0]
    plot.has_data_labels = chart_spec.get("show_values", kind != "line")
    if plot.has_data_labels:
        plot.data_labels.show_value = True
        plot.data_labels.font.size = Pt(8)
        if chart_spec.get("number_format"):
            plot.data_labels.number_format = chart_spec["number_format"]
            plot.data_labels.number_format_is_linked = False
    if hasattr(plot, "gap_width"):
        plot.gap_width = 85
    chart.value_axis.major_gridlines.format.line.color.rgb = generate.RULE
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
    chart.value_axis.format.line.fill.background()
    chart.category_axis.format.line.fill.background()
    for series, color in zip(
            chart.series,
            (generate.GRAY, generate.ACCENT, generate.CORAL)):
        if kind == "line":
            series.format.line.color.rgb = color
            series.format.line.width = Pt(1.6)
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = color
            series.marker.format.line.color.rgb = color
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color


def _render_diagram(slide, side, child, area):
    row_count = len(child["diagram"]["rows"])
    natural_h = min(area.height, 1.65 + row_count * 0.72)
    top_area = generate.ContentArea(
        area.top, area.top + natural_h, area.shifted,
        area.left, area.right)
    render_diagram(slide, child["diagram"], content_area=top_area)


_CHILD_RENDERERS = {
    "bullets": _render_bullets,
    "cards": _render_cards,
    "image": _render_image,
    "table": _render_table,
    "chart": _render_chart,
    "diagram": _render_diagram,
}


def s_split(slide, spec, page):
    """共通ヘッダー下を左右へ分け、対応する子レイアウタを独立描画する。"""
    area = generate.header(
        slide, spec["kicker"], spec["title"], spec.get("lead"))
    regions = _regions(area, bool(spec.get("note")))
    right = regions[1]
    generate.add_rect(
        slide, right.left, right.top - 0.02, right.width,
        right.height - 0.02, generate.ZEBRA)
    for side, child, region in zip(
            ("left", "right"), (spec["left"], spec["right"]), regions):
        child = dict(child, _region=region)
        content_area = _content_region(slide, side, child)
        _CHILD_RENDERERS[child["type"]](slide, side, child, content_area)
    if spec.get("note"):
        generate.note_line(slide, spec["note"])
