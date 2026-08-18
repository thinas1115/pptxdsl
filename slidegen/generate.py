"""python-pptxとテキスト実測による基本サンプルデッキ生成。"""
import argparse
import math
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from content import DECK
from cover_footer import load_cover_footer_config, render_cover, render_footer
from layout_fit import (
    FitError,
    fit_text_or_raise,
    fit_vertical_stacks,
    select_fit,
    stepped,
)
from textfit import line_height_in, text_width_in, wrap_natural, wrap_text

# ---- テーマ ----
NAVY = RGBColor(0x18, 0x2C, 0x43)
ACCENT = RGBColor(0x0D, 0x78, 0x70)
CORAL = RGBColor(0xC7, 0x58, 0x3E)
LIGHT = RGBColor(0xDF, 0xEB, 0xE8)
TEXT = RGBColor(0x20, 0x27, 0x29)
GRAY = RGBColor(0x66, 0x6E, 0x70)
DONE_TEXT = RGBColor(0x56, 0x60, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFC)
SURFACE = RGBColor(0xFA, 0xF9, 0xF5)
ZEBRA = RGBColor(0xEC, 0xEA, 0xE4)
CANVAS = RGBColor(0xF7, 0xF5, 0xEF)
RULE = RGBColor(0xD1, 0xCF, 0xC8)
FONT = "Yu Gothic"
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.55
BODY_W = SLIDE_W - MARGIN * 2
BODY_TOP, BODY_BOTTOM = 1.58, 6.85
LEAD_Y, LEAD_MAX_H = 1.58, 0.56
TABLE_HEADER_H = 0.72
TABLE_TOP_GAP = 0.38
TABLE_BOTTOM_GAP = 0.15
TABLE_NOTE_H = 0.30
COVER_FOOTER = load_cover_footer_config()


@dataclass(frozen=True)
class ContentArea:
    """ヘッダー下のrenderer描画領域。leadなしでは従来値を保持する。"""

    top: float = BODY_TOP
    bottom: float = BODY_BOTTOM
    shifted: bool = False

    @property
    def height(self):
        return self.bottom - self.top

    def map_y(self, y):
        """固定構図の従来Y座標を、lead指定時の本文領域へ写像する。"""
        if not self.shifted:
            return y
        scale = self.height / (BODY_BOTTOM - BODY_TOP)
        return self.top + (y - BODY_TOP) * scale


def configure_cover_footer(path=None):
    """表紙・フッター設定を切り替える。未指定なら標準設定へ戻す。"""
    global COVER_FOOTER
    COVER_FOOTER = load_cover_footer_config(path)


def set_run(run, size, *, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "ja-JP")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def add_text(slide, x, y, w, h, text, size, *, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.3,
             wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        set_run(p.add_run(), size, bold=bold, color=color)
        p.runs[0].text = line
    return tb


def add_rect(slide, x, y, w, h, fill, *, line=None, round_=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        sp.adjustments[0] = 0.06
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def header(slide, kicker, title, lead=None):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CANVAS)
    kicker_size, _ = fit_text_or_raise(
        "header", "kicker", kicker, 4.8, 0.32, 11.5,
        min_pt=9, weight="bold", spacing=1.1)
    add_text(slide, 0.72, 0.27, 4.8, 0.32, kicker, kicker_size,
             bold=True, color=ACCENT)
    size = 27
    lines = wrap_natural(title, 11.9, size, "bold")
    while len(lines) > 1 and size > 18:
        size -= 0.5
        lines = wrap_natural(title, 11.9, size, "bold")
    if len(lines) > 1:
        size, lines = fit_text_or_raise(
            "header", "title", title, 11.9, 0.86, 18,
            min_pt=16, weight="bold", spacing=1.12)
    add_text(slide, 0.72, 0.67, 11.9, 0.86, "\n".join(lines), size,
             bold=True, color=NAVY, spacing=1.12)
    if not lead:
        return ContentArea()

    lead_size, lead_lines = fit_text_or_raise(
        "header", "lead", lead, 11.9, LEAD_MAX_H, 14,
        min_pt=11.5, spacing=1.18)
    lead_h = len(lead_lines) * line_height_in(lead_size, 1.18) + 0.03
    add_text(slide, 0.72, LEAD_Y, 11.9, lead_h, "\n".join(lead_lines),
             lead_size, color=GRAY, spacing=1.18)
    return ContentArea(top=LEAD_Y + lead_h + 0.17, shifted=True)


def page_label(page):
    """Return a stable page marker shared by every generator entry point."""
    total = len(DECK["slides"])
    digits = max(2, len(str(total)))
    return f"{page:0{digits}d} / {total:0{digits}d}"


def footer(slide, page):
    total = len(DECK["slides"])
    render_footer(slide, page, DECK["meta"], total, COVER_FOOTER,
                  add_text=add_text, add_rect=add_rect)


def note_line(slide, note, link=None):
    link_label = link["label"] if link else ""
    display_text = f"{note} {link_label}" if link else note
    size, _ = fit_text_or_raise(
        "note", "text", display_text, BODY_W, 0.25, 8.5,
        min_pt=7, spacing=1.1)
    tb = add_text(slide, MARGIN, 6.62, BODY_W, 0.25, note, size,
                  color=GRAY, align=PP_ALIGN.RIGHT)
    if link:
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        set_run(run, size, color=ACCENT)
        run.font.underline = True
        run.text = f" {link_label}"
        run.hyperlink.address = link["url"]


# ---- スライド種別 ----
def s_title(slide, spec, page):
    meta = DECK["meta"]
    total = len(DECK["slides"])
    render_cover(slide, spec, meta, total, COVER_FOOTER,
                 add_text=add_text, add_rect=add_rect)


def _normalize_bullet(item):
    """公開object形式と既存の2要素配列を同じ描画入力へ揃える。"""
    if isinstance(item, dict):
        return item["text"], bool(item.get("checked", False))
    return item[0], False


def _bullet_first_line_center(y, size):
    """Return the visual center of the first body-text line in inches."""
    return y + size / 108 + 0.012


def s_bullets(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    style = spec.get("style", "numbered")
    bullets = [_normalize_bullet(item) for item in spec["bullets"]]
    top_gap = 0.27 if area.shifted else 0.32
    area_h = area.height - top_gap - 0.28
    tx, tw = 1.48, 9.65
    preferred_size = 22 if len(bullets) <= 2 else 20 if len(bullets) <= 4 else 18

    fitted = fit_vertical_stacks(
        "bullets", area_h, [bullets],
        lambda item, size: (
            len(wrap_text(item[0], tw, size))
            * line_height_in(size, 1.22)
        ),
        standard_size=preferred_size,
        min_size=12,
        font_step=0.5,
        standard_gap=0.52,
        min_gap=0.34,
        gap_step=0.03,
        guidance="箇条書きを減らすか各項目を短くしてください。",
    )
    size, gap = fitted.size, fitted.gap
    heights = fitted.stacks[0]
    rule_w = min(9.05, max(
        6.30,
        max(text_width_in(text, size) for text, _checked in bullets) + 0.30,
    ))
    y = area.top + top_gap
    for i, ((text, checked), bh) in enumerate(zip(bullets, heights), 1):
        marker_center = _bullet_first_line_center(y, size)
        if style == "numbered":
            # Latin numerals sit slightly below the geometric center when
            # vertically anchored, so keep the text box 0.012in higher.
            add_text(slide, 0.80, marker_center - 0.162, 0.46, 0.30,
                     f"{i:02d}", 15.5, bold=True, color=ACCENT,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        elif style == "bullet":
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1.04), Inches(marker_center - 0.07),
                Inches(0.14), Inches(0.14))
            marker.fill.solid()
            marker.fill.fore_color.rgb = ACCENT
            marker.line.fill.background()
            marker.shadow.inherit = False
        else:
            checkbox_top = marker_center - 0.12
            add_rect(
                slide, 0.99, checkbox_top, 0.24, 0.24,
                ACCENT if checked else CANVAS,
                line=ACCENT if checked else GRAY,
            )
            if checked:
                add_text(slide, 0.99, checkbox_top - 0.028, 0.24, 0.25, "✓", 12,
                         bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE, spacing=1.0, wrap=False)
        add_text(slide, tx, y, tw, bh + 0.08, text, size,
                 color=DONE_TEXT if style == "checklist" and checked else TEXT,
                 spacing=1.22)
        if style == "numbered":
            add_rect(slide, tx, y + bh + 0.16, rule_w, 0.01, RULE)
        y += bh + gap


def s_cards(slide, spec, page):
    """Render purpose-specific cards instead of one universal panel pattern."""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    cards = [_normalize_card(card) for card in spec["cards"]]
    n = len(cards)
    style = spec.get("style", "editorial")
    left, usable_w = 0.78, 11.78

    if style == "metrics":
        cols = n if n <= 4 else 3
        rows = math.ceil(n / cols)
        gap_x, gap_y = 0.46, 0.42
        cw = (usable_w - gap_x * (cols - 1)) / cols
        ch = min(2.72, (area.height - 0.58 - gap_y * (rows - 1)) / rows)
        if ch < 1.62:
            raise FitError(
                "cards.metrics: KPIカードの高さが不足しています。カード数または本文を"
                "減らしてください。")
        top = area.top + 0.42
        body_size = min(
            fit_text_or_raise(
                "cards.metrics", f"cards[{i}].body", card["body"],
                cw - 0.16, ch - 1.26, 13, min_pt=10.5, spacing=1.16,
            )[0]
            for i, card in enumerate(cards))
        for i, card in enumerate(cards):
            row, col = divmod(i, cols)
            x = left + col * (cw + gap_x)
            y = top + row * (ch + gap_y)
            label, value = card["heading"], card.get("value", "")
            if not value:
                label, value = _split_metric_head(label)
            label_size, _ = fit_text_or_raise(
                "cards.metrics", f"cards[{i}].label", label,
                cw, 0.34, 12.5, min_pt=10.5, weight="bold", spacing=1.1)
            add_text(slide, x, y, cw, 0.34, label, label_size,
                     bold=True, color=NAVY)
            value_size = fit_text_or_raise(
                "cards.metrics", f"cards[{i}].value", value,
                cw, 0.66, 32, min_pt=22, weight="bold")[0]
            color = CORAL if card["emphasis"] else ACCENT
            add_text(slide, x, y + 0.42, cw, 0.66, value, value_size,
                     bold=True, color=color)
            add_text(slide, x, y + 1.22, cw, ch - 1.22, card["body"],
                     body_size, color=TEXT, spacing=1.2)
            if col < cols - 1 and i + 1 < n:
                add_rect(slide, x + cw + gap_x / 2, y + 0.04,
                         0.012, ch - 0.08, RULE)
            if row < rows - 1:
                add_rect(slide, x, y + ch + gap_y / 2, cw, 0.012, RULE)
        return

    emphasized = [i for i, card in enumerate(cards) if card["emphasis"]]
    if style == "editorial" and n == 4 and len(emphasized) == 1:
        lead_index = emphasized[0]
        lead_card = cards[lead_index]
        supporting = [card for i, card in enumerate(cards) if i != lead_index]
        top = area.top + 0.46
        add_text(slide, 0.8, top, 0.58, 0.36, "01", 15,
                 bold=True, color=GRAY)
        lead_head_size, _ = fit_text_or_raise(
            "cards.editorial", f"cards[{lead_index}].heading", lead_card["heading"],
            3.95, 0.56, 23, min_pt=18, weight="bold", spacing=1.1)
        add_text(slide, 1.52, top - 0.03, 3.95, 0.56, lead_card["heading"],
                 lead_head_size, bold=True, color=NAVY)
        lead_size, lead_lines = fit_text_or_raise(
            "cards.editorial", f"cards[{lead_index}].body", lead_card["body"],
            4.1, 2.05, 16, min_pt=13, spacing=1.28)
        add_text(slide, 1.52, top + 0.76, 4.1, 2.05,
                 "\n".join(lead_lines), lead_size, color=TEXT, spacing=1.28)

        right_x, right_w = 6.42, 5.92
        row_h = 1.31
        for i, card in enumerate(supporting, 2):
            y = top + (i - 2) * row_h
            add_text(slide, right_x, y + 0.02, 0.5, 0.32, f"{i:02d}", 12.5,
                     bold=True, color=GRAY)
            head_size, _ = fit_text_or_raise(
                "cards.editorial", f"supporting[{i - 2}].heading", card["heading"],
                right_w - 0.68, 0.38, 16, min_pt=13,
                weight="bold", spacing=1.1)
            add_text(slide, right_x + 0.68, y, right_w - 0.68, 0.38,
                     card["heading"],
                     head_size, bold=True, color=NAVY)
            body_size, body_lines = fit_text_or_raise(
                "cards.editorial", f"supporting[{i - 2}].body", card["body"],
                right_w - 0.72, 0.56, 12.5, min_pt=11, spacing=1.15)
            add_text(slide, right_x + 0.68, y + 0.48, right_w - 0.72, 0.56,
                     "\n".join(body_lines), body_size,
                     color=TEXT, spacing=1.15)
            if i < n:
                add_rect(slide, right_x + 0.68, y + 1.14, right_w - 0.68, 0.012, RULE)
        return

    cols = n if n <= 3 else (2 if n == 4 else 3)
    rows = math.ceil(n / cols)
    gap_x, gap_y = 0.72, 0.44
    cw = (usable_w - gap_x * (cols - 1)) / cols
    area_h = area.height - 0.62
    fit_available = area_h if rows == 2 else min(3.32, area_h)

    row_cards = [
        cards[row * cols:(row + 1) * cols]
        for row in range(rows)
    ]
    fitted = fit_vertical_stacks(
        "cards", fit_available, [row_cards],
        lambda cards_in_row, size: (
            0.72 + max(
                len(wrap_text(card["body"], cw - 0.92, size))
                * line_height_in(size, 1.2)
                for card in cards_in_row
            )
        ),
        standard_size=14,
        min_size=11,
        font_step=0.5,
        standard_gap=gap_y,
        min_gap=0.28 if rows == 2 else gap_y,
        gap_step=0.04,
        guidance="カード本文を短くするかカード数を減らしてください。",
    )
    body_size = fitted.size
    gap_y = fitted.gap
    ch = ((area_h - gap_y * (rows - 1)) / rows
          if rows > 1 else min(3.32, area_h))
    top = area.top + (0.42 if rows == 2 else 0.72)
    for i, card in enumerate(cards):
        row, col = divmod(i, cols)
        x = left + col * (cw + gap_x)
        y = top + row * (ch + gap_y)
        add_text(slide, x, y + 0.02, 0.58, 0.36, f"{i + 1:02d}", 14,
                 bold=True, color=GRAY)
        head_size, _ = fit_text_or_raise(
            "cards", f"cards[{i}].heading", card["heading"],
            cw - 0.72, 0.42, 16.5, min_pt=13,
            weight="bold", spacing=1.1)
        head_color = CORAL if card["emphasis"] else NAVY
        add_text(slide, x + 0.72, y, cw - 0.72, 0.42,
                 card["heading"], head_size, bold=True, color=head_color)
        add_text(slide, x + 0.72, y + 0.58, cw - 0.82, ch - 0.64, card["body"],
                 body_size, color=TEXT, spacing=1.2)
        if row < rows - 1:
            add_rect(slide, x + 0.72, y + ch + gap_y / 2, cw - 0.72, 0.012, RULE)


def _split_metric_head(head):
    match = re.match(r"^(.*?)[\s　]+([+\-−]?[0-9][0-9,.]*\s*(?:%|分|件|倍|pt)?)$", head)
    if match:
        return match.group(1).strip(), match.group(2).replace("−", "-").strip()
    return head, ""


def _normalize_card(card):
    """旧2要素配列と新しい意味付きobjectを共通形式へ変換する。"""
    if isinstance(card, dict):
        return {
            "heading": card["heading"],
            "body": card["body"],
            "value": card.get("value", ""),
            "emphasis": bool(card.get("emphasis")),
        }
    return {
        "heading": card[0], "body": card[1], "value": "", "emphasis": False,
    }


def _auto_table_widths(columns, rows):
    """見出しとセル内容を実測し、本文幅を列へ決定論的に配分する。"""
    count = len(columns)
    min_width = max(0.92, min(1.38, BODY_W / count * 0.58))
    natural = []
    for index, heading in enumerate(columns):
        samples = [text_width_in(heading, 12, "bold")]
        samples.extend(text_width_in(row[index], 11.5) for row in rows)
        natural.append(min(4.65, max(min_width, max(samples) + 0.34)))

    widths = [min_width] * count
    remaining = BODY_W - sum(widths)
    unmet = {index for index in range(count) if natural[index] > widths[index]}
    while remaining > 0.001 and unmet:
        share = remaining / len(unmet)
        used = 0.0
        completed = set()
        for index in unmet:
            addition = min(share, natural[index] - widths[index])
            widths[index] += addition
            used += addition
            if natural[index] - widths[index] <= 0.001:
                completed.add(index)
        remaining -= used
        unmet -= completed
        if used <= 0.001:
            break
    if remaining > 0:
        weights = [max(1.0, value) for value in natural]
        total_weight = sum(weights)
        widths = [width + remaining * weight / total_weight
                  for width, weight in zip(widths, weights)]
    widths[-1] += BODY_W - sum(widths)
    return widths


def _fit_table(rows, widths, avail):
    """表の余白圧縮・フォント縮小を選び、行高と判定結果を返す。"""
    min_row_h = min(1.04, max(0.64, avail / max(1, len(rows))))

    def measure(size, pad):
        row_hs = []
        for row in rows:
            need = max(
                len(wrap_text(c, widths[j] - pad * 2, size))
                * line_height_in(size, 1.15) for j, c in enumerate(row))
            row_hs.append(max(need + pad * 2, min_row_h))
        return row_hs, sum(row_hs)

    def candidates():
        for pad in stepped(0.14, 0.10, 0.02):
            _row_hs, used = measure(13.5, pad)
            yield ("standard" if pad == 0.14 else "padding",
                   {"size": 13.5, "pad": pad}, used)
        for size in stepped(13.0, 10.5, 0.5):
            _row_hs, used = measure(size, 0.10)
            yield "font", {"size": size, "pad": 0.10}, used

    fitted = select_fit(
        "table", avail, candidates(),
        guidance="表の行を減らすかセル内の文言を短くしてください。",
    )
    row_hs, _used = measure(fitted.values["size"], fitted.values["pad"])
    return fitted, row_hs


def s_table(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    cols, rows = spec["columns"], spec["rows"]
    widths = spec.get("col_widths") or _auto_table_widths(cols, rows)
    assert abs(sum(widths) - BODY_W) < 0.6, f"列幅合計={sum(widths)}"
    hdr_h = TABLE_HEADER_H
    table_available = (
        area.height - TABLE_TOP_GAP - TABLE_BOTTOM_GAP
        - (TABLE_NOTE_H if spec.get("note") else 0)
    )
    avail = table_available - hdr_h
    fitted, row_hs = _fit_table(rows, widths, avail)
    size, pad = fitted.values["size"], fitted.values["pad"]
    table_h = hdr_h + sum(row_hs)
    top = area.top + TABLE_TOP_GAP + max(0.0, (table_available - table_h) * 0.12)
    gt = slide.shapes.add_table(len(rows) + 1, len(cols), Inches(MARGIN),
                                Inches(top), Inches(BODY_W), Inches(table_h))
    table = gt.table
    table.first_row = False
    table.horz_banding = False
    for j, wdt in enumerate(widths):
        table.columns[j].width = Emu(int(Inches(wdt)))
    table.rows[0].height = Emu(int(Inches(hdr_h)))
    for i, rh in enumerate(row_hs):
        table.rows[i + 1].height = Emu(int(Inches(rh)))
    for j, name in enumerate(cols):
        header_size, _ = fit_text_or_raise(
            "table", f"columns[{j}]", name, widths[j] - 0.18, hdr_h - 0.08,
            size, min_pt=10.5, weight="bold", spacing=1.15)
        _cell(table.cell(0, j), name, header_size, bold=True, color=WHITE, fill=NAVY,
              center=False)
    for i, row in enumerate(rows):
        fill = SURFACE if i % 2 else ZEBRA
        for j, val in enumerate(row):
            _cell(table.cell(i + 1, j), val, size,
                  bold=(j == 0), color=NAVY if j == 0 else TEXT, fill=fill,
                  center=False)
    if spec.get("note"):
        note_line(slide, spec["note"], spec.get("note_link"))


def _cell(cell, text, size, *, bold=False, color=TEXT, fill=SURFACE, center=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.09)
    cell.margin_top = cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        p.line_spacing = 1.15
        set_run(p.add_run(), size, bold=bold, color=color)
        p.runs[0].text = line


def s_twocol(slide, spec, page):
    """Render a purposeful before/after comparison, not generic cards."""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    gap = 0.42
    left = 0.76
    cw = (11.82 - gap) / 2
    panels = [spec["left"], spec["right"]]
    text_w = cw - 1.10

    fitted = fit_vertical_stacks(
        "two_column", area.height - 0.42,
        [panel["bullets"] for panel in panels],
        lambda bullet, size: max(
            0.42,
            len(wrap_text(bullet, text_w, size))
            * line_height_in(size, 1.18) + 0.06,
        ),
        standard_size=14,
        min_size=11,
        font_step=0.5,
        standard_gap=0.18,
        min_gap=0.10,
        gap_step=0.02,
        fixed_height=1.63,
        guidance="左右の箇条書きを減らすか文言を短くしてください。",
    )
    size, row_gap = fitted.size, fitted.gap
    rows = fitted.stacks
    panel_h = max(2.72, fitted.used)
    top = area.top + 0.30
    for i, p in enumerate(panels):
        x = left + i * (cw + gap)
        fill = ZEBRA if i == 0 else LIGHT
        marker = GRAY if i == 0 else ACCENT
        add_rect(slide, x, top, cw, panel_h, fill)
        add_text(slide, x + 0.38, top + 0.29, cw - 0.76, 0.22,
                 p.get("label", "BEFORE" if i == 0 else "AFTER"), 9.5,
                 bold=True, color=marker)
        heading_size, _ = fit_text_or_raise(
            "two_column", f"{'left' if i == 0 else 'right'}.heading",
            p["heading"], cw - 0.76, 0.42, 18,
            min_pt=14, weight="bold", spacing=1.1)
        add_text(slide, x + 0.38, top + 0.64, cw - 0.76, 0.42,
                 p["heading"], heading_size, bold=True, color=NAVY)
        add_rect(slide, x + 0.38, top + 1.16, cw - 0.76, 0.01, RULE)
        y = top + 1.34
        for row_index, (bullet, row_h) in enumerate(zip(p["bullets"], rows[i])):
            add_text(slide, x + 0.38, y + 0.01, 0.34, 0.24,
                     f"{row_index + 1:02d}", 9.5, bold=True, color=marker)
            add_text(slide, x + 0.82, y, text_w, row_h, bullet, size,
                     color=TEXT, spacing=1.18)
            y += row_h
            if row_index < len(rows[i]) - 1:
                add_rect(slide, x + 0.82, y + row_gap / 2,
                         text_w, 0.008, RULE)
                y += row_gap


def s_chart(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    categories = spec["chart"]["categories"]
    series = spec["chart"]["series"]
    kind = spec["chart"].get("kind", "bar")
    if kind not in {"bar", "column", "line", "stacked_bar", "stacked_column"}:
        raise FitError(f"chart.kind={kind!r} は未対応です。")
    if not 1 <= len(categories) <= 12 or not 1 <= len(series) <= 4:
        raise FitError(
            "chart: 描画可能範囲はカテゴリ1〜12件、系列1〜4件です。"
            "カテゴリまたは系列を減らしてください。")
    _s_chart_modern(slide, spec, area, categories, series, kind)
    if spec.get("note"):
        note_line(slide, spec["note"])


def _s_chart_modern(slide, spec, area, categories, series, kind):
    show_legend = spec["chart"].get("show_legend", len(series) > 1)
    bounds = _modern_chart_bounds(area, show_legend=show_legend)
    palette = _modern_chart_palette(len(series))
    if kind == "bar":
        _modern_bar_chart(slide, bounds, categories, series, palette, stacked=False)
    elif kind == "stacked_bar":
        _modern_bar_chart(slide, bounds, categories, series, palette, stacked=True)
    elif kind == "column":
        _modern_column_chart(slide, bounds, categories, series, palette, stacked=False)
    elif kind == "stacked_column":
        _modern_column_chart(slide, bounds, categories, series, palette, stacked=True)
    elif kind == "line":
        _modern_line_chart(slide, bounds, categories, series, palette)
    if show_legend:
        _modern_legend(slide, bounds, series, palette)


def _modern_chart_palette(count):
    return [ACCENT, NAVY, CORAL, GRAY][:count]


def _modern_chart_bounds(area, *, show_legend):
    legend_h = 0.42 if show_legend else 0.0
    return {
        "x": 0.95,
        "y": area.top + 0.25,
        "w": 11.55,
        "h": max(2.65, area.height - 0.62 - legend_h),
        "legend_h": legend_h,
    }


def _chart_line(slide, x1, y1, x2, y2, color=RULE, width=0.8):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def _chart_dot(slide, x, y, color, size=0.08, *, line_color=None):
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x - size / 2), Inches(y - size / 2),
        Inches(size), Inches(size))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.color.rgb = line_color or color
    dot.line.width = Pt(0.6)
    return dot


def _format_chart_value(value):
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    if isinstance(value, float):
        return f"{value:.1f}"
    return str(value)


def _modern_axis_scale(values):
    vmax = max([0.0] + [float(v) for v in values])
    if vmax <= 0:
        return 1.0, [0, 0.25, 0.5, 0.75, 1.0]
    raw = vmax / 4
    magnitude = 10 ** math.floor(math.log10(raw))
    step = next(s * magnitude for s in (1, 2, 5, 10) if s * magnitude >= raw)
    axis_max = step * 4
    return axis_max, [step * i for i in range(5)]


def _all_values(series, *, stacked, category_count):
    if stacked:
        return [
            sum(float(vals[i]) for _, vals in series)
            for i in range(category_count)
        ]
    return [float(v) for _, vals in series for v in vals]


def _modern_value_axis(slide, plot_x, plot_y, plot_w, plot_h, axis_max, ticks):
    for tick in ticks:
        ratio = 0 if axis_max == 0 else tick / axis_max
        y = plot_y + plot_h - plot_h * ratio
        _chart_line(slide, plot_x, y, plot_x + plot_w, y, RULE, 0.55)
        add_text(slide, plot_x - 0.46, y - 0.09, 0.36, 0.18,
                 _format_chart_value(tick), 7.8, color=GRAY, align=PP_ALIGN.RIGHT)
    _chart_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, GRAY, 0.8)
    _chart_line(slide, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h, GRAY, 0.8)


def _modern_horizontal_axis(slide, plot_x, plot_y, plot_w, plot_h, axis_max, ticks):
    for tick in ticks:
        ratio = 0 if axis_max == 0 else tick / axis_max
        x = plot_x + plot_w * ratio
        _chart_line(slide, x, plot_y, x, plot_y + plot_h, RULE, 0.55)
        add_text(slide, x - 0.18, plot_y + plot_h + 0.08, 0.36, 0.18,
                 _format_chart_value(tick), 7.8, color=GRAY, align=PP_ALIGN.CENTER)
    _chart_line(slide, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h, GRAY, 0.8)


def _modern_bar_chart(slide, bounds, categories, series, palette, *, stacked):
    cat_count = len(categories)
    values = _all_values(series, stacked=stacked, category_count=cat_count)
    axis_max, ticks = _modern_axis_scale(values)
    label_w = 1.52
    plot_x = bounds["x"] + label_w + 0.20
    plot_y = bounds["y"] + 0.12
    plot_w = bounds["w"] - label_w - 0.35
    plot_h = bounds["h"] - 0.18
    row_h = plot_h / cat_count
    _modern_horizontal_axis(slide, plot_x, plot_y, plot_w, plot_h, axis_max, ticks)
    for i, cat in enumerate(categories):
        y_mid = plot_y + row_h * i + row_h / 2
        add_text(slide, bounds["x"], y_mid - 0.15, label_w, 0.30,
                 str(cat), 9.4, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
        if stacked:
            x_cursor = plot_x
            for sidx, (_, vals) in enumerate(series):
                val = float(vals[i])
                bw = plot_w * val / axis_max if axis_max else 0
                add_rect(slide, x_cursor, y_mid - 0.13, max(0.02, bw), 0.26, palette[sidx])
                if bw > 0.42:
                    add_text(slide, x_cursor + 0.04, y_mid - 0.10, bw - 0.08, 0.20,
                             _format_chart_value(val), 7.8, bold=True, color=WHITE,
                             align=PP_ALIGN.CENTER)
                x_cursor += bw
        else:
            bar_gap = 0.04
            bar_h = min(0.22, (row_h - 0.16) / len(series) - bar_gap)
            group_y = y_mid - (bar_h * len(series) + bar_gap * (len(series) - 1)) / 2
            for sidx, (_, vals) in enumerate(series):
                val = float(vals[i])
                y = group_y + sidx * (bar_h + bar_gap)
                bw = plot_w * val / axis_max if axis_max else 0
                add_rect(slide, plot_x, y, max(0.02, bw), bar_h, palette[sidx])
                add_text(slide, plot_x + bw + 0.06, y - 0.02, 0.56, bar_h + 0.04,
                         _format_chart_value(val), 7.6, color=GRAY)


def _modern_column_chart(slide, bounds, categories, series, palette, *, stacked):
    cat_count = len(categories)
    values = _all_values(series, stacked=stacked, category_count=cat_count)
    axis_max, ticks = _modern_axis_scale(values)
    plot_x = bounds["x"] + 0.60
    plot_y = bounds["y"] + 0.12
    plot_w = bounds["w"] - 0.78
    plot_h = bounds["h"] - 0.56
    _modern_value_axis(slide, plot_x, plot_y, plot_w, plot_h, axis_max, ticks)
    cluster_w = plot_w / cat_count
    for i, cat in enumerate(categories):
        x0 = plot_x + cluster_w * i
        add_text(slide, x0 + 0.03, plot_y + plot_h + 0.10, cluster_w - 0.06, 0.26,
                 str(cat), 8.2, color=GRAY, align=PP_ALIGN.CENTER)
        if stacked:
            bar_w = min(0.58, cluster_w * 0.44)
            x = x0 + (cluster_w - bar_w) / 2
            y_base = plot_y + plot_h
            for sidx, (_, vals) in enumerate(series):
                val = float(vals[i])
                h = plot_h * val / axis_max if axis_max else 0
                add_rect(slide, x, y_base - h, bar_w, max(0.02, h), palette[sidx])
                y_base -= h
        else:
            gap = 0.035
            bar_w = min(0.28, (cluster_w * 0.64 - gap * (len(series) - 1)) / len(series))
            group_w = bar_w * len(series) + gap * (len(series) - 1)
            for sidx, (_, vals) in enumerate(series):
                val = float(vals[i])
                h = plot_h * val / axis_max if axis_max else 0
                x = x0 + (cluster_w - group_w) / 2 + sidx * (bar_w + gap)
                add_rect(slide, x, plot_y + plot_h - h, bar_w, max(0.02, h), palette[sidx])


def _modern_line_chart(slide, bounds, categories, series, palette):
    values = _all_values(series, stacked=False, category_count=len(categories))
    axis_max, ticks = _modern_axis_scale(values)
    plot_x = bounds["x"] + 0.60
    plot_y = bounds["y"] + 0.12
    plot_w = bounds["w"] - 0.78
    plot_h = bounds["h"] - 0.56
    _modern_value_axis(slide, plot_x, plot_y, plot_w, plot_h, axis_max, ticks)
    step = plot_w / max(1, len(categories) - 1)
    for i, cat in enumerate(categories):
        x = plot_x + step * i if len(categories) > 1 else plot_x + plot_w / 2
        add_text(slide, x - 0.34, plot_y + plot_h + 0.10, 0.68, 0.24,
                 str(cat), 8.0, color=GRAY, align=PP_ALIGN.CENTER)
    for sidx, (_, vals) in enumerate(series):
        points = []
        for i, val in enumerate(vals):
            x = plot_x + step * i if len(categories) > 1 else plot_x + plot_w / 2
            y = plot_y + plot_h - plot_h * float(val) / axis_max
            points.append((x, y))
        for (x1, y1), (x2, y2) in zip(points, points[1:]):
            _chart_line(slide, x1, y1, x2, y2, palette[sidx], 1.8)
        for x, y in points:
            _chart_dot(slide, x, y, palette[sidx], 0.10, line_color=WHITE)


def _modern_legend(slide, bounds, series, palette):
    y = bounds["y"] + bounds["h"] + 0.20
    x = bounds["x"] + 0.62
    for idx, (name, _) in enumerate(series):
        lx = x + idx * 2.05
        add_rect(slide, lx, y + 0.07, 0.12, 0.12, palette[idx])
        add_text(slide, lx + 0.18, y, 1.72, 0.26, name, 8.2, color=GRAY)


RENDER = {"title": s_title, "bullets": s_bullets, "cards": s_cards,
          "table": s_table, "two_column": s_twocol, "chart": s_chart}


def render_slide(renderer, slide, spec, idx):
    """rendererの収容エラーをスライド位置つきの運用メッセージへ変換する。"""
    try:
        renderer(slide, spec, idx)
    except (ValueError, FileNotFoundError) as e:
        raise SystemExit(
            f"NG: slides[{idx - 1}] (type={spec['type']}) の生成に失敗:\n"
            f"  {e}") from e


def main(out_path, cover_footer_config=None):
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        configure_cover_footer(cover_footer_config)
    except ValueError as e:
        raise SystemExit(f"NG: 表紙・フッター設定: {e}") from e
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    total = len(DECK["slides"])
    for idx, spec in enumerate(DECK["slides"], 1):
        slide = prs.slides.add_slide(blank)
        render_slide(RENDER[spec["type"]], slide, spec, idx)
        if spec["type"] != "title":
            footer(slide, idx)
    prs.save(out_path)
    print(f"saved: {out_path} ({total} slides)")


if __name__ == "__main__":
    default_out = Path(__file__).resolve().parent.parent / "out" / "sample_basic.pptx"
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("out_path", nargs="?", default=default_out)
    parser.add_argument("--cover-footer-config", metavar="PATH",
                        help="表紙・フッター設定JSON")
    args = parser.parse_args()
    main(args.out_path, args.cover_footer_config)
