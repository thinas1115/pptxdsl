"""生成済みPPTXのレイアウト衝突を機械検知する品質ゲート。

使い方: python slidegen/check_layout.py out/deck.pptx
検知対象:
  T-T: テキストグリフ同士の交差
  T-P: テキストグリフ×画像(アイコン)
  T-S: テキストグリフ×塗り図形の「部分重なり」(完全内包=意図的デザインは許可)
  T-F: テキストグリフがコンテナ枠線をまたぐ
  L-T: 矢印・線×テキストグリフ(白塗りマスクラベルは除外)
  CELL-OOB: 表セルからのテキストはみ出し
  OOB: スライド境界からの図形・画像・表・グラフのはみ出し
  VIS-CONTRAST: 背景に接する意味面とキャンバスの分離不足
  SEQ-CLEARANCE: sequence自己処理の戻り線とメッセージラベルの接触
"""
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from cover_footer import COVER_BACKGROUND_NAME
from quality_markers import (
    MIN_SURFACE_CONTRAST,
    MIN_SURFACE_EDGE_CONTRAST,
    SEQUENCE_LABEL_CLEARANCE,
    SEQUENCE_MESSAGE_LABEL_PREFIX,
    SEQUENCE_SELF_ROUTE_PREFIX,
    SURFACE_ON_CANVAS_PREFIX,
)
from textfit import line_height_in, wrap_text

EMU = 914400
EDGE = 0.03          # 枠線の当たり判定幅
SEG_TRIM = 0.08      # 線分端はノード接続なので判定から除外する長さ
EPS = 0.03           # 視認できない接触(スリバー)を無視する許容量
BACKGROUND_PREFIX = "layout-background:"


def rect_of(sh):
    return (sh.left / EMU, sh.top / EMU, (sh.left + sh.width) / EMU,
            (sh.top + sh.height) / EMU)


def intersects(a, b, eps=EPS):
    return not (a[2] - eps <= b[0] or b[2] - eps <= a[0]
                or a[3] - eps <= b[1] or b[3] - eps <= a[1])


def expanded(rect, amount):
    return (rect[0] - amount, rect[1] - amount,
            rect[2] + amount, rect[3] + amount)


def contains(outer, inner, eps=EPS):
    return (outer[0] - eps <= inner[0] and outer[1] - eps <= inner[1]
            and outer[2] + eps >= inner[2] and outer[3] + eps >= inner[3])


def glyph_rect_for_text_frame(tf, box):
    """指定領域にあるtext frameの実グリフ矩形をtextfitで推定する。"""
    bw = box[2] - box[0]
    lines, sizes = [], []
    align = PP_ALIGN.LEFT
    for p in tf.paragraphs:
        if not p.runs:
            continue
        size = p.runs[0].font.size.pt if p.runs[0].font.size else 11.0
        bold = bool(p.runs[0].font.bold)
        text = "".join(r.text for r in p.runs)
        weight = "bold" if bold else "regular"
        for ln in wrap_text(text, bw, size, weight):
            lines.append((ln, size, weight))
            sizes.append(size)
        if p.alignment is not None:
            align = p.alignment
    if not lines:
        return None
    from textfit import text_width_in
    w = max(text_width_in(t, s, wt) for t, s, wt in lines)
    h = sum(line_height_in(s) for _, s, _ in lines)
    if align == PP_ALIGN.CENTER:
        x0 = box[0] + (bw - w) / 2
    elif align == PP_ALIGN.RIGHT:
        x0 = box[2] - w
    else:
        x0 = box[0]
    if tf.vertical_anchor == MSO_ANCHOR.MIDDLE:
        y0 = box[1] + ((box[3] - box[1]) - h) / 2
    else:
        y0 = box[1]
    return (x0, y0, x0 + w, y0 + h)


def glyph_rect(sh):
    """通常図形の実グリフ矩形をtextfitで推定する。"""
    return glyph_rect_for_text_frame(sh.text_frame, rect_of(sh))


def table_cells(sh):
    """表セルの外形・文字領域とtext frameを返す。"""
    table = sh.table
    x0, y0, _, _ = rect_of(sh)
    col_widths = [col.width / EMU for col in table.columns]
    row_heights = [row.height / EMU for row in table.rows]
    y = y0
    for ri, row in enumerate(table.rows):
        x = x0
        row_h = row_heights[ri]
        for ci, col in enumerate(table.columns):
            col_w = col_widths[ci]
            cell = table.cell(ri, ci)
            if not getattr(cell, "is_spanned", False):
                span_w = getattr(cell, "span_width", 1)
                span_h = getattr(cell, "span_height", 1)
                merged_w = sum(col_widths[ci:ci + span_w])
                merged_h = sum(row_heights[ri:ri + span_h])
                outer = (x, y, x + merged_w, y + merged_h)
                inner = (
                    x + cell.margin_left / EMU,
                    y + cell.margin_top / EMU,
                    x + merged_w - cell.margin_right / EMU,
                    y + merged_h - cell.margin_bottom / EMU,
                )
                yield outer, inner, cell.text_frame, ri, ci
            x += col_w
        y += row_h


def is_oob(rect, slide_w, slide_h):
    return (rect[0] < -EPS or rect[1] < -EPS
            or rect[2] > slide_w + EPS or rect[3] > slide_h + EPS)


def seg_of(sh):
    """コネクタの線分端点(flipで向きを解決)。"""
    x1, y1, x2, y2 = rect_of(sh)
    fh = getattr(sh, "rotation", 0)  # dummy no-op
    el = sh._element
    flip_h = el.spPr.xfrm.get("flipH") == "1" if el.spPr.xfrm is not None else False
    flip_v = el.spPr.xfrm.get("flipV") == "1" if el.spPr.xfrm is not None else False
    ax, bx = (x2, x1) if flip_h else (x1, x2)
    ay, by = (y2, y1) if flip_v else (y1, y2)
    return (ax, ay, bx, by)


def seg_hits_rect(seg, r, eps=EPS):
    """線分と矩形の交差(端をSEG_TRIMだけ縮めて接続点は無視)。"""
    x1, y1, x2, y2 = seg
    import math
    L = math.hypot(x2 - x1, y2 - y1)
    if L < SEG_TRIM * 2:
        return False
    ux, uy = (x2 - x1) / L, (y2 - y1) / L
    x1, y1 = x1 + ux * SEG_TRIM, y1 + uy * SEG_TRIM
    x2, y2 = x2 - ux * SEG_TRIM, y2 - uy * SEG_TRIM
    # 双方の投影が重なるかをサンプリングで簡易判定
    steps = max(2, int(L / 0.05))
    for i in range(steps + 1):
        t = i / steps
        px, py = x1 + (x2 - x1) * t, y1 + (y2 - y1) * t
        if r[0] + eps < px < r[2] - eps and r[1] + eps < py < r[3] - eps:
            return True
    return False


def has_solid_fill(sh):
    try:
        return sh.fill.type is not None and str(sh.fill.type) == "SOLID (1)"
    except Exception:
        return False


def solid_fill_rgb(sh):
    if not has_solid_fill(sh):
        return None
    try:
        return tuple(sh.fill.fore_color.rgb)
    except (AttributeError, TypeError):
        return None


def solid_line_rgb(sh):
    """図形の単色輪郭をRGBで返す。輪郭なし・テーマ色は判定不能とする。"""
    try:
        rgb = sh.line.color.rgb
        if rgb is None or sh.line.width <= 0:
            return None
        return tuple(rgb)
    except (AttributeError, TypeError, ValueError):
        return None


def relative_luminance(rgb):
    values = [value / 255 for value in rgb]
    linear = [
        value / 12.92 if value <= 0.04045
        else ((value + 0.055) / 1.055) ** 2.4
        for value in values
    ]
    return 0.2126 * linear[0] + 0.7152 * linear[1] + 0.0722 * linear[2]


def contrast_ratio(first, second):
    high, low = sorted(
        (relative_luminance(first), relative_luminance(second)), reverse=True)
    return (high + 0.05) / (low + 0.05)


def canvas_fill_rgb(slide, slide_w, slide_h):
    """全画面を覆う最背面の塗り図形をキャンバスとして扱う。"""
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        bounds = rect_of(sh)
        if (bounds[0] <= EPS and bounds[1] <= EPS
                and bounds[2] >= slide_w - EPS
                and bounds[3] >= slide_h - EPS):
            color = solid_fill_rgb(sh)
            if color is not None:
                return color
    return None


def snippet(t):
    t = t.replace("\n", " ")
    return t[:16] + ("…" if len(t) > 16 else "")


def check(path):
    prs = Presentation(path)
    slide_w = prs.slide_width / EMU
    slide_h = prs.slide_height / EMU
    findings = []
    for si, slide in enumerate(prs.slides, 1):
        canvas_rgb = canvas_fill_rgb(slide, slide_w, slide_h)
        texts, pics, solids, frames, segs = [], [], [], [], []
        sequence_returns, sequence_labels = [], []
        for z, sh in enumerate(slide.shapes):  # zは描画順(後勝ち)
            st = sh.shape_type
            bounds = rect_of(sh)
            if is_oob(bounds, slide_w, slide_h):
                findings.append((si, "OOB", sh.name, ""))
            if st == MSO_SHAPE_TYPE.PICTURE:
                if sh.name != COVER_BACKGROUND_NAME:
                    pics.append((bounds, sh.name))
            elif st in (MSO_SHAPE_TYPE.AUTO_SHAPE,):
                if sh.name.startswith(SURFACE_ON_CANVAS_PREFIX):
                    surface_rgb = solid_fill_rgb(sh)
                    if surface_rgb is None or canvas_rgb is None:
                        findings.append((si, "VIS-CONTRAST", sh.name,
                                         "塗りまたはキャンバス色を取得できません"))
                    else:
                        fill_ratio = contrast_ratio(surface_rgb, canvas_rgb)
                        edge_rgb = solid_line_rgb(sh)
                        edge_ratio = (
                            contrast_ratio(edge_rgb, canvas_rgb)
                            if edge_rgb is not None else 0.0
                        )
                        if (fill_ratio < MIN_SURFACE_CONTRAST
                                and edge_ratio < MIN_SURFACE_EDGE_CONTRAST):
                            findings.append((
                                si, "VIS-CONTRAST", sh.name,
                                f"面={fill_ratio:.3f} < {MIN_SURFACE_CONTRAST:.2f}, "
                                f"輪郭={edge_ratio:.3f} < "
                                f"{MIN_SURFACE_EDGE_CONTRAST:.2f}"))
                if not sh.name.startswith(BACKGROUND_PREFIX):
                    (solids if has_solid_fill(sh) else frames).append(
                        (bounds, sh.name, z))
            elif st == MSO_SHAPE_TYPE.LINE:
                seg = seg_of(sh)
                segs.append((seg, sh.name, z))
                if (sh.name.startswith(SEQUENCE_SELF_ROUTE_PREFIX)
                        and sh.name.endswith(":return")):
                    sequence_returns.append((seg, sh.name))
            elif st == MSO_SHAPE_TYPE.CHART:
                pics.append((bounds, sh.name))
            elif st == MSO_SHAPE_TYPE.TABLE:
                for _cell_rect, text_rect, tf, ri, ci in table_cells(sh):
                    if not tf.text.strip():
                        continue
                    g = glyph_rect_for_text_frame(tf, text_rect)
                    if g:
                        label = snippet(tf.text)
                        texts.append((g, label, False))
                        if not contains(text_rect, g, eps=0.01):
                            findings.append(
                                (si, "CELL-OOB", f"{sh.name}[{ri},{ci}]", label))
            if st != MSO_SHAPE_TYPE.TABLE \
                    and sh.has_text_frame and sh.text_frame.text.strip():
                g = glyph_rect(sh)
                if g:
                    texts.append((g, snippet(sh.text_frame.text),
                                  has_solid_fill(sh)))
                if sh.name.startswith(SEQUENCE_MESSAGE_LABEL_PREFIX):
                    sequence_labels.append((bounds, sh.name))
        # T-T
        for i in range(len(texts)):
            for j in range(i + 1, len(texts)):
                if intersects(texts[i][0], texts[j][0]):
                    findings.append((si, "T-T", texts[i][1], texts[j][1]))
        # T-P
        for g, t, _ in texts:
            for r, name in pics:
                if intersects(g, r):
                    findings.append((si, "T-P", t, name))
        # T-S: 部分重なりのみ(内包は許可)
        for g, t, _ in texts:
            for r, name, _z in solids:
                if intersects(g, r) and not contains(r, g):
                    findings.append((si, "T-S", t, name))
        # T-F: 枠線またぎ(内包/完全外は許可。白塗りマスクラベルは枠線を隠すので許可)
        for g, t, masked in texts:
            if masked:
                continue
            for r, name, _z in frames:
                if intersects(g, r, eps=-EDGE) and not contains(
                        (r[0] + EDGE, r[1] + EDGE, r[2] - EDGE, r[3] - EDGE), g) \
                        and intersects(g, r):
                    findings.append((si, "T-F", t, name))
        # L-T (マスクラベルと、線より後に描かれた塗り図形上のテキストは除外)
        for seg, name, lz in segs:
            for g, t, masked in texts:
                if masked or not seg_hits_rect(seg, g):
                    continue
                covered = any(contains(r, g) and sz > lz
                              for r, _n, sz in solids)
                if not covered:
                    findings.append((si, "L-T", name, t))
        # L-P
        for seg, name, _lz in segs:
            for r, pname in pics:
                if seg_hits_rect(seg, r):
                    findings.append((si, "L-P", name, pname))
        # 線上ラベルは塗りマスクを使うため汎用L-Tでは除外される。
        # 自己処理の戻り線は、別メッセージのラベルに隠れると経路が欠けるため意味名で検査する。
        for seg, route_name in sequence_returns:
            for label_rect, label_name in sequence_labels:
                if seg_hits_rect(
                        seg, expanded(label_rect, SEQUENCE_LABEL_CLEARANCE),
                        eps=0.0):
                    findings.append((
                        si, "SEQ-CLEARANCE", route_name, label_name))
        for g, t, _ in texts:
            if is_oob(g, slide_w, slide_h):
                findings.append((si, "OOB-TEXT", t, ""))
    return findings


if __name__ == "__main__":
    path = sys.argv[1]
    fs = check(path)
    if not fs:
        print(f"OK: no layout or visual QA findings in {path}")
        sys.exit(0)
    print(f"NG: {len(fs)} finding(s) in {path}")
    for si, kind, a, b in fs:
        print(f"  slide{si:02d} [{kind}] {a!r} x {b!r}")
    sys.exit(1)
