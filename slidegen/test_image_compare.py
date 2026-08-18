"""image_compare rendererのパス・段階的収容・schemaを検証する。"""
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

import generate
from image_slide import (
    COMPARE_PICTURE_NAME, fit_image_compare_layout, s_image_compare,
)
from layout_fit import FitError
from validate_content import validate


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    return prs.slides.add_slide(prs.slide_layouts[6])


def _must_fail(fn, expected):
    try:
        fn()
    except FitError as exc:
        assert expected in str(exc), str(exc)
    else:
        raise AssertionError("過密入力を拒否しませんでした")


def main():
    spec = {
        "type": "image_compare", "kicker": "比較", "title": "画像比較",
        "left": {"image": "images/pptxdsl-repository.png", "label": "変更前"},
        "right": {"image": "cover/cover-background.png", "label": "変更後"},
        "fit": "cover", "shadow": True,
    }
    deck = {
        "meta": {"title": "検証", "footer": "検証", "date": "2026年7月",
                 "author": "検証担当"},
        "slides": [spec],
    }
    assert not validate(deck)

    slide = _slide()
    s_image_compare(slide, spec, 1)
    pictures = [shape for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 2
    assert all(p.name == COMPARE_PICTURE_NAME for p in pictures)
    left_pic, right_pic = sorted(pictures, key=lambda p: p.left)
    assert left_pic.left < right_pic.left
    # cover指定なので、縦横比に応じてどちらかの軸をトリミングしている。
    assert (left_pic.crop_top > 0 or left_pic.crop_left > 0)
    assert (right_pic.crop_top > 0 or right_pic.crop_left > 0)

    texts = {shape.text_frame.text for shape in slide.shapes
             if shape.has_text_frame}
    assert {"変更前", "変更後"} <= texts

    contain_spec = dict(spec, fit="contain")
    contain_slide = _slide()
    s_image_compare(contain_slide, contain_spec, 1)
    contain_pictures = [shape for shape in contain_slide.shapes
                         if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert all(p.crop_left == p.crop_right == 0
               and p.crop_top == p.crop_bottom == 0
               for p in contain_pictures)

    default_label_spec = deepcopy(spec)
    del default_label_spec["left"]["label"]
    del default_label_spec["right"]["label"]
    default_slide = _slide()
    s_image_compare(default_slide, default_label_spec, 1)
    default_texts = {shape.text_frame.text for shape in default_slide.shapes
                      if shape.has_text_frame}
    assert {"BEFORE", "AFTER"} <= default_texts

    assert fit_image_compare_layout(5.27).stage == "standard"
    assert fit_image_compare_layout(3.60).stage == "gap"
    assert fit_image_compare_layout(3.20).stage == "element"
    _must_fail(lambda: fit_image_compare_layout(2.50), "最小設定")

    invalid = deepcopy(deck)
    del invalid["slides"][0]["right"]
    assert any('"right"' in error for error in validate(invalid))

    invalid = deepcopy(deck)
    invalid["slides"][0]["left"]["image"] = "images/missing.png"
    assert any("left.image" in error and "assets/にありません" in error
               for error in validate(invalid))

    invalid = deepcopy(deck)
    invalid["slides"][0]["right"]["image"] = "../../outside.png"
    assert any("assets内" in error for error in validate(invalid))

    invalid = deepcopy(deck)
    invalid["slides"][0]["fit"] = "stretch"
    assert any("contain" in error for error in validate(invalid))

    invalid = deepcopy(deck)
    invalid["slides"][0]["shadow"] = "yes"
    assert any("true または false" in error for error in validate(invalid))

    invalid = deepcopy(deck)
    invalid["slides"][0]["left"]["caption"] = "説明"
    assert any("left" in error and "caption" in error for error in validate(invalid))

    print("image_compare slide tests passed")


if __name__ == "__main__":
    main()
