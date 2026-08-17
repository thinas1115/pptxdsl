"""check_layoutが壊れたPPTXを誤って合格させないための回帰テスト。"""
from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

from check_layout import check
from quality_markers import (
    SEQUENCE_MESSAGE_LABEL_PREFIX,
    SEQUENCE_SELF_ROUTE_PREFIX,
    SURFACE_ON_CANVAS_PREFIX,
)


def save(prs, path):
    prs.save(path)
    return check(path)


with TemporaryDirectory() as td:
    out = Path(td)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_table(1, 2, Inches(1), Inches(7.1), Inches(8), Inches(1))
    findings = save(prs, out / "table_oob.pptx")
    assert any(kind == "OOB" for _, kind, _, _ in findings), findings

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = ChartData()
    data.categories = ["A", "B"]
    data.add_series("値", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(12.8), Inches(1),
        Inches(1), Inches(2), data)
    findings = save(prs, out / "chart_oob.pptx")
    assert any(kind == "OOB" for _, kind, _, _ in findings), findings

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(9.8), Inches(1), Inches(0.5), Inches(0.5))
    findings = save(prs, out / "custom_size_oob.pptx")
    assert any(kind == "OOB" for _, kind, _, _ in findings), findings

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide.shapes.add_table(
        1, 1, Inches(1), Inches(1), Inches(1.2), Inches(0.25)).table
    cell = table.cell(0, 0)
    cell.text = "セルから確実にはみ出す長い文章です"
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    findings = save(prs, out / "cell_text_oob.pptx")
    assert any(kind == "CELL-OOB" for _, kind, _, _ in findings), findings

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(5), Inches(1.5)).table
    merged = table.cell(0, 0)
    merged.merge(table.cell(0, 1))
    merged.text = "結合セルの幅全体を使う見出し"
    findings = save(prs, out / "merged_table_ok.pptx")
    assert not any(kind == "CELL-OOB" for _, kind, _, _ in findings), findings

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    canvas = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    canvas.fill.solid()
    canvas.fill.fore_color.rgb = RGBColor(0xF7, 0xF5, 0xEF)
    stage = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(1), Inches(1), Inches(3), Inches(0.5))
    stage.name = f"{SURFACE_ON_CANVAS_PREFIX}test-stage"
    stage.fill.solid()
    stage.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFC)
    stage.line.fill.background()
    findings = save(prs, out / "surface_contrast_ng.pptx")
    assert any(kind == "VIS-CONTRAST" for _, kind, _, _ in findings), findings

    stage.line.color.rgb = RGBColor(0xD1, 0xCF, 0xC8)
    stage.line.width = Pt(0.85)
    findings = save(prs, out / "surface_outline_contrast_ok.pptx")
    assert not any(kind == "VIS-CONTRAST" for _, kind, _, _ in findings), findings

    stage.line.fill.background()
    stage.fill.fore_color.rgb = RGBColor(0xDF, 0xEB, 0xE8)
    findings = save(prs, out / "surface_fill_contrast_ok.pptx")
    assert not any(kind == "VIS-CONTRAST" for _, kind, _, _ in findings), findings

    # 塗りマスク付きラベルが自己処理の戻り線を隠す回帰を、汎用L-Tとは別に拒否する。
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    returned = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(2), Inches(4), Inches(2))
    returned.name = f"{SEQUENCE_SELF_ROUTE_PREFIX}self:return"
    label = slide.shapes.add_textbox(
        Inches(2), Inches(1.90), Inches(1.2), Inches(0.20))
    label.name = f"{SEQUENCE_MESSAGE_LABEL_PREFIX}next"
    label.text_frame.text = "次の処理"
    label.fill.solid()
    label.fill.fore_color.rgb = RGBColor(0xF7, 0xF5, 0xEF)
    findings = save(prs, out / "sequence_label_masks_return_ng.pptx")
    assert any(kind == "SEQ-CLEARANCE" for _, kind, _, _ in findings), findings

    label.top = Inches(2.02)
    findings = save(prs, out / "sequence_label_clearance_ng.pptx")
    assert any(kind == "SEQ-CLEARANCE" for _, kind, _, _ in findings), findings

    label.top = Inches(2.05)
    findings = save(prs, out / "sequence_label_clearance_ok.pptx")
    assert not any(kind == "SEQ-CLEARANCE" for _, kind, _, _ in findings), findings

print("check_layout broken-PPTX regression: ALL OK")
