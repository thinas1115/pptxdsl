"""技術研修rendererの標準入力、過密入力、停止条件、意味検証を確認する。"""

import json
from copy import deepcopy
from pathlib import Path
from types import SimpleNamespace

from pptx import Presentation
from pptx.util import Inches

import generate
from layout_fit import FitError
from training_renderers import (
    _fit_protocol_state_flow,
    s_code_lab,
    s_concept,
    s_knowledge_check,
    s_network,
    s_protocol_state_flow,
    s_protocol_anatomy,
)
from validate_content import validate


RENDERERS = {
    "concept": s_concept,
    "nw_topology": s_network,
    "nw_protocol_flow": s_protocol_state_flow,
    "nw_frame_anatomy": s_protocol_anatomy,
    "config_lab": s_code_lab,
    "knowledge_check": s_knowledge_check,
}


def _presentation():
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    return prs


def _render(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    RENDERERS[spec["type"]](slide, spec, len(prs.slides))
    return slide


def _assert_in_slide(slide):
    tolerance = Inches(0.01)
    for shape in slide.shapes:
        assert shape.left >= -tolerance
        assert shape.top >= -tolerance
        assert shape.left + shape.width <= Inches(generate.SLIDE_W) + tolerance
        assert shape.top + shape.height <= Inches(generate.SLIDE_H) + tolerance


def _samples():
    path = Path(__file__).resolve().parent.parent / "examples" / "training" / "vlan" / "content.json"
    deck = json.loads(path.read_text(encoding="utf-8"))
    result = {}
    for spec in deck["slides"]:
        if spec["type"] in RENDERERS and spec["type"] not in result:
            result[spec["type"]] = spec
    return deck, result


def _assert_network_contract(deck, samples):
    network = next(
        deepcopy(spec) for spec in deck["slides"]
        if (spec["type"] == "nw_topology" and len(spec["lanes"]) > 1
            and any(link.get("kind") == "access" for link in spec["links"]))
    )
    access_index = next(
        index for index, link in enumerate(network["links"])
        if link.get("kind") == "access"
    )
    access_link = network["links"][access_index]
    source = next(node for node in network["nodes"] if node["id"] == access_link["from"])
    invalid_lane = next(
        lane["id"] for lane in network["lanes"]
        if lane["id"] not in source["lanes"]
    )
    access_link["lanes"] = [invalid_lane]
    errors = validate({"meta": deck["meta"], "slides": [network]},
                      allow_sample_content=True)
    assert any("双方に所属するlane" in error for error in errors), errors

    missing_access_lane = deepcopy(network)
    missing_access_lane["links"][access_index]["lanes"] = []
    errors = validate(
        {"meta": deck["meta"], "slides": [missing_access_lane]},
        allow_sample_content=True,
    )
    assert any("accessのlanesは1件" in error for error in errors), errors

    trunk = next(
        deepcopy(spec) for spec in deck["slides"]
        if spec["type"] == "nw_topology"
        and any(link.get("kind") == "trunk" for link in spec["links"])
    )

    slide = _render(_presentation(), trunk)
    texts = {
        shape.text.strip() for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    }
    assert any("VLAN 10 / 20" in text for text in texts)
    assert any(shape.name.startswith("layout-background:network-lane:")
               for shape in slide.shapes)

    paired = deepcopy(trunk)
    source = next(node for node in paired["nodes"] if len(node["lanes"]) == 1)
    duplicate = deepcopy(source)
    duplicate["id"] = "paired-device"
    duplicate["label"] = "予備端末"
    paired["nodes"].append(duplicate)
    errors = validate(
        {"meta": deck["meta"], "slides": [paired]},
        allow_sample_content=True,
    )
    assert not errors, errors
    paired_slide = _render(_presentation(), paired)
    labels = {
        shape.text: shape.left
        for shape in paired_slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text in {source["label"], duplicate["label"]}
    }
    assert len(labels) == 2
    assert labels[source["label"]] != labels[duplicate["label"]]

    overcrowded = deepcopy(paired)
    third = deepcopy(source)
    third["id"] = "third-device"
    third["label"] = "3台目"
    overcrowded["nodes"].append(third)
    errors = validate(
        {"meta": deck["meta"], "slides": [overcrowded]},
        allow_sample_content=True,
    )
    assert any("2件まで配置" in error for error in errors), errors

    repeated_link = deepcopy(trunk)
    reverse = deepcopy(repeated_link["links"][0])
    reverse["from"], reverse["to"] = reverse["to"], reverse["from"]
    repeated_link["links"].append(reverse)
    errors = validate(
        {"meta": deck["meta"], "slides": [repeated_link]},
        allow_sample_content=True,
    )
    assert any("向きを問わず重複" in error for error in errors), errors

    blocked = next(
        deepcopy(spec) for spec in deck["slides"]
        if spec["type"] == "nw_topology"
        and any(link.get("kind") == "blocked" for link in spec["links"])
    )
    blocked_index = next(
        index for index, link in enumerate(blocked["links"])
        if link.get("kind") == "blocked"
    )
    invalid_blocked = deepcopy(blocked)
    invalid_blocked["links"][blocked_index]["lanes"] = ["v20"]
    errors = validate(
        {"meta": deck["meta"], "slides": [invalid_blocked]},
        allow_sample_content=True,
    )
    assert any("接続元だけが所属" in error for error in errors), errors
    blocked_slide = _render(_presentation(), blocked)
    assert any(
        getattr(shape, "has_text_frame", False) and shape.text == "×"
        for shape in blocked_slide.shapes
    )


def _assert_concept_contract(samples):
    concept = deepcopy(samples["concept"])
    rendered = _render(_presentation(), concept)
    texts = {
        shape.text.strip() for shape in rendered.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    }
    assert concept["term"] in texts
    assert concept["misconception"] in texts

    dense = deepcopy(concept)
    dense["definition"] = "復旧目標の起点と終点を合意する説明です。" * 40
    try:
        _render(_presentation(), dense)
    except FitError:
        pass
    else:
        raise AssertionError("過密な定義文を明示停止できていません")


def _assert_protocol_state_flow_contract(samples):
    protocol_state_flow = deepcopy(samples["nw_protocol_flow"])
    rendered = _render(_presentation(), protocol_state_flow)
    texts = {
        shape.text.strip() for shape in rendered.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    }
    assert "802.1Qタグ付き" in texts
    assert "Accessポートで分類" in texts
    assert any(
        shape.name.startswith(
            "layout-background:protocol-state-flow:track:"
        )
        for shape in rendered.shapes
    )

    appearances = deepcopy(protocol_state_flow)
    for state, appearance in zip(
            appearances["flows"][0]["states"],
            ("plain", "encapsulated", "internal", "alert", "plain")):
        state["appearance"] = appearance
        if appearance == "encapsulated":
            state["encapsulation"] = "HDR"
        else:
            state.pop("encapsulation", None)
    appearance_slide = _render(_presentation(), appearances)
    marker_names = {shape.name for shape in appearance_slide.shapes}
    for appearance in ("plain", "encapsulated", "internal", "alert"):
        assert f"protocol-state-flow:marker:{appearance}" in marker_names

    minimum = deepcopy(protocol_state_flow)
    minimum["stages"] = minimum["stages"][:3]
    minimum["flows"] = [minimum["flows"][0]]
    minimum["flows"][0]["states"] = minimum["flows"][0]["states"][:3]
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [minimum]},
        allow_sample_content=True,
    )
    assert not errors, errors
    _assert_in_slide(_render(_presentation(), minimum))

    maximum = deepcopy(protocol_state_flow)
    maximum.pop("lead", None)
    maximum["title"] = "最大構成"
    maximum["stages"].append({
        "id": "archive", "label": "記録先",
        "icon": "icons/fluent/database.png", "role": "endpoint",
    })
    for flow in maximum["flows"]:
        flow.pop("sub", None)
        for state in flow["states"]:
            state.pop("detail", None)
        flow["states"].append({
            "stage": "archive", "label": "記録", "appearance": "plain",
        })
    third_flow = deepcopy(maximum["flows"][0])
    third_flow["label"] = "処理結果"
    for state in third_flow["states"]:
        state["label"] = "正常"
    maximum["flows"].append(third_flow)
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [maximum]},
        allow_sample_content=True,
    )
    assert not errors, errors
    _assert_in_slide(_render(_presentation(), maximum))

    def observed_stages(spec, heights):
        result_stages = set()
        for height in heights:
            try:
                result = _fit_protocol_state_flow(
                    SimpleNamespace(height=height),
                    spec["stages"], spec["flows"],
                )
            except FitError:
                continue
            result_stages.add(result.stage)
        return result_stages

    observed_fit_stages = observed_stages(
        protocol_state_flow,
        (5.0 - index * 0.01 for index in range(180)),
    )
    detail_stress = deepcopy(protocol_state_flow)
    for flow in detail_stress["flows"]:
        for state in flow["states"]:
            state["detail"] = (state.get("detail", "補足") + "判断条件") * 2
    observed_fit_stages.update(observed_stages(
        detail_stress,
        (4.2 - index * 0.01 for index in range(90)),
    ))
    icon_stress = deepcopy(maximum)
    icon_stress["flows"] = [icon_stress["flows"][0]]
    icon_stress["lead"] = "追跡対象の状態変化を段階ごとに確認する。"
    for stage in icon_stress["stages"]:
        stage["label"] = "非常に長い処理段階の表示名称と補足情報"
    observed_fit_stages.update(observed_stages(
        icon_stress,
        (3.2 - index * 0.01 for index in range(140)),
    ))
    assert {"standard", "gap", "icon", "font"} <= observed_fit_stages, (
        observed_fit_stages
    )

    missing = deepcopy(protocol_state_flow)
    missing["flows"][0]["states"].pop()
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [missing]},
        allow_sample_content=True,
    )
    assert any("全stageを1件ずつ" in error for error in errors), errors

    duplicate = deepcopy(protocol_state_flow)
    duplicate["flows"][0]["states"][-1]["stage"] = "tx"
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [duplicate]},
        allow_sample_content=True,
    )
    assert any("が重複しています" in error for error in errors), errors

    duplicate_flow = deepcopy(protocol_state_flow)
    duplicate_flow["flows"][1]["label"] = duplicate_flow["flows"][0]["label"]
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [duplicate_flow]},
        allow_sample_content=True,
    )
    assert any("flows[1].label" in error and "重複" in error
               for error in errors), errors

    invalid_role = deepcopy(protocol_state_flow)
    invalid_role["stages"][0]["role"] = "gateway"
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [invalid_role]},
        allow_sample_content=True,
    )
    assert any("endpoint / processor / link" in error for error in errors), errors

    invalid_appearance = deepcopy(protocol_state_flow)
    invalid_appearance["flows"][0]["states"][0]["appearance"] = "tagged"
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [invalid_appearance]},
        allow_sample_content=True,
    )
    assert any("plain / encapsulated / internal / alert" in error
               for error in errors), errors

    missing_encapsulation = deepcopy(protocol_state_flow)
    state = missing_encapsulation["flows"][0]["states"][0]
    state["appearance"] = "encapsulated"
    state.pop("encapsulation", None)
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [missing_encapsulation]},
        allow_sample_content=True,
    )
    assert any("encapsulation はencapsulatedで必須" in error
               for error in errors), errors

    invalid_encapsulation = deepcopy(protocol_state_flow)
    invalid_encapsulation["flows"][0]["states"][0]["encapsulation"] = "TAG"
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [invalid_encapsulation]},
        allow_sample_content=True,
    )
    assert any("appearance=encapsulatedでのみ" in error
               for error in errors), errors

    long_encapsulation = deepcopy(protocol_state_flow)
    state = long_encapsulation["flows"][0]["states"][0]
    state["appearance"] = "encapsulated"
    state["encapsulation"] = "TOO-LONG-HDR"
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [long_encapsulation]},
        allow_sample_content=True,
    )
    assert any("encapsulation は8文字以内" in error
               for error in errors), errors

    unknown_nested_key = deepcopy(protocol_state_flow)
    unknown_nested_key["flows"][0]["states"][0]["x"] = 0.4
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [unknown_nested_key]},
        allow_sample_content=True,
    )
    assert any("states[0].x" in error and "未対応" in error
               for error in errors), errors

    dense = deepcopy(protocol_state_flow)
    dense.pop("lead", None)
    dense["flows"] = []
    for flow_index in range(3):
        flow = deepcopy(protocol_state_flow["flows"][flow_index % 2])
        flow["label"] = f"長い比較対象 {flow_index + 1}"
        for state in flow["states"]:
            state["label"] = "変換後の状態を詳細に説明する長い見出し"
            state["detail"] = "処理内容と判断条件を省略せずに記載した長い補足説明" * 3
        dense["flows"].append(flow)
    try:
        _render(_presentation(), dense)
    except FitError:
        pass
    else:
        raise AssertionError("過密なprotocol_state_flowを明示停止できていません")


def _assert_protocol_contract(samples):
    protocol = deepcopy(samples["nw_frame_anatomy"])
    protocol["frames"][0]["annotations"] = [
        {"field": "missing", "text": "未定義フィールド"}
    ]
    errors = validate(
        {"meta": {"title": "検証"}, "slides": [protocol]},
        allow_sample_content=True)
    assert any("定義済みfield" in error for error in errors), errors

    variable = deepcopy(samples["nw_frame_anatomy"])
    variable["frames"][0]["fields"][-1]["size_label"] = "可変長"
    rendered = _render(_presentation(), variable)
    texts = {
        shape.text.strip() for shape in rendered.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    }
    assert "TPID" in texts
    assert "可変長" in texts
    assert "tpid" not in texts

    comparison = deepcopy(samples["nw_frame_anatomy"])
    base_frame = comparison["frames"][0]
    tagged_frame = deepcopy(base_frame)
    base_frame["label"] = "タグなし"
    tagged_frame["label"] = "タグ付き"
    tagged_frame["fields"].insert(
        -1, {"id": "tag", "name": "VLANタグ", "bits": 32,
             "role": "highlight"})
    comparison["frames"] = [base_frame, tagged_frame]
    rendered = _render(_presentation(), comparison)
    comparison_texts = {
        shape.text.strip() for shape in rendered.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    }
    assert "基準" in comparison_texts
    assert "+4 byte" in comparison_texts
    baseline_label = next(
        shape for shape in rendered.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text.strip() == "基準")
    tagged_label = next(
        shape for shape in rendered.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text.strip() == "+4 byte")
    assert tagged_label.left > baseline_label.left
    common_fields = [
        shape for shape in rendered.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text.strip() == "送信元MAC"
    ]
    assert len(common_fields) == 2
    assert common_fields[0].width == common_fields[1].width

    dense = deepcopy(samples["nw_frame_anatomy"])
    dense.pop("lead", None)
    template = dense["frames"][0]
    annotations = [
        {"field": field["id"], "text": f"{field['name']}の役割"}
        for field in template["fields"][:4]
    ]
    dense["frames"] = []
    for index in range(3):
        frame = deepcopy(template)
        frame["label"] = f"Frame {index + 1}"
        frame["annotations"] = deepcopy(annotations)
        dense["frames"].append(frame)
    rendered = _render(_presentation(), dense)
    _assert_in_slide(rendered)


def _assert_code_stop(samples):
    code = deepcopy(samples["config_lab"])
    code["sections"][0]["code"] = "x" * 240
    try:
        _render(_presentation(), code)
    except FitError:
        pass
    else:
        raise AssertionError("折り返せない長いコードを明示停止できていません")


def _assert_knowledge_modes(deck):
    questions = next(
        deepcopy(spec) for spec in deck["slides"]
        if spec["type"] == "knowledge_check" and spec["mode"] == "questions")
    answers = next(
        deepcopy(spec) for spec in deck["slides"]
        if spec["type"] == "knowledge_check" and spec["mode"] == "answers")
    question_texts = {
        shape.text.strip() for shape in _render(_presentation(), questions).shapes
        if getattr(shape, "has_text_frame", False)
    }
    answer_texts = {
        shape.text.strip() for shape in _render(_presentation(), answers).shapes
        if getattr(shape, "has_text_frame", False)
    }
    assert "解説" not in question_texts
    assert "解説" in answer_texts

    long_options = deepcopy(questions)
    long_options["questions"] = long_options["questions"][:1]
    long_options["questions"][0]["options"][0] = (
        "同一ブロードキャストドメインの範囲と転送経路を確認してから設定を変更する"
    )
    _assert_in_slide(_render(_presentation(), long_options))


def main():
    deck, samples = _samples()
    errors = validate(deepcopy(deck), allow_sample_content=True)
    assert not errors, "\n".join(errors)
    assert set(samples) == set(RENDERERS)

    prs = _presentation()
    for spec in samples.values():
        _assert_in_slide(_render(prs, deepcopy(spec)))

    for type_ in ("nw_protocol_flow", "nw_frame_anatomy", "config_lab"):
        legacy = deepcopy(samples[type_])
        legacy["takeaway"] = "下部注記帯へ表示していた要点"
        legacy_errors = validate(
            {"meta": {"title": "検証"}, "slides": [legacy]},
            allow_sample_content=True,
        )
        assert any(
            '"takeaway"' in error and '"lead"' in error
            for error in legacy_errors
        ), legacy_errors

    _assert_network_contract(deck, samples)
    _assert_concept_contract(samples)
    _assert_protocol_state_flow_contract(samples)
    _assert_protocol_contract(samples)
    _assert_code_stop(samples)
    _assert_knowledge_modes(deck)
    print("training renderer tests: OK")


if __name__ == "__main__":
    main()
