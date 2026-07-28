"""左右分割rendererのschema、描画、過密停止を検証する。"""
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches

import generate
from layout_fit import FitError
from split_layout import s_split
from validate_content import validate


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    return prs.slides.add_slide(prs.slide_layouts[6])


def _diagram():
    return {
        "cols": ["entry", "service", "data"],
        "rows": ["main"],
        "nodes": {
            "user": {
                "col": "entry", "row": "main",
                "icon": "icons/fluent/person.png", "title": "利用者",
            },
            "api": {
                "col": "service", "row": "main",
                "icon": "icons/fluent/globe.png", "title": "業務API",
            },
            "db": {
                "col": "data", "row": "main",
                "icon": "icons/fluent/database.png", "title": "データベース",
            },
        },
        "containers": [],
        "channels": {},
        "edges": [
            {"from": "user", "to": "api", "label": "HTTPS"},
            {"from": "api", "to": "db", "label": "SQL"},
        ],
    }


def _spec(left, right, lead=None):
    spec = {
        "type": "split",
        "kicker": "複合レイアウト",
        "title": "異なる情報構造を左右で対応付ける",
        "left": left,
        "right": right,
    }
    if lead:
        spec["lead"] = lead
    return spec


def _deck(spec):
    return {"meta": {"title": "左右分割検証"}, "slides": [spec]}


def _must_fit_fail(spec):
    try:
        s_split(_slide(), spec, 1)
    except FitError:
        return
    raise AssertionError("半幅の最小設定でも収まらない入力を拒否しませんでした")


def main():
    diagram_cards = _spec(
        {"type": "diagram", "heading": "処理経路", "diagram": _diagram()},
        {
            "type": "cards", "heading": "設計上の要点",
            "cards": [
                {"heading": "境界", "body": "公開領域とデータ領域を分離する。"},
                {"heading": "監視", "body": "主要な接続点で状態を確認する。"},
                {"heading": "復旧", "body": "障害時の切替手順を事前に定める。"},
            ],
        },
    )
    image_bullets = _spec(
        {
            "type": "image", "heading": "対象画面",
            "image": "images/pptxdsl-repository.png", "fit": "contain",
        },
        {
            "type": "bullets", "heading": "確認事項",
            "bullets": [
                ["操作の起点を確認する", None],
                ["主要な導線を追跡する", None],
                ["権限境界を明確にする", None],
            ],
        },
        lead="画面と確認事項を同じ視野で照合します。",
    )
    chart_table = _spec(
        {
            "type": "chart", "heading": "月別推移",
            "chart": {
                "kind": "line",
                "categories": ["4月", "5月", "6月", "7月"],
                "series": [["実績", [12, 18, 27, 35]]],
            },
        },
        {
            "type": "table", "heading": "判断基準",
            "columns": ["指標", "基準"],
            "rows": [["品質", "95%以上"], ["時間", "8時間以内"]],
        },
    )
    bullets_cards = _spec(
        {
            "type": "bullets", "heading": "前提",
            "bullets": [["対象部門を限定する", None], ["評価期間を定める", None]],
        },
        {
            "type": "cards", "heading": "判断材料",
            "cards": [
                {"heading": "品質", "body": "正確性を確認する。"},
                {"heading": "運用", "body": "継続負荷を確認する。"},
            ],
        },
    )
    for spec in (diagram_cards, image_bullets, chart_table, bullets_cards):
        errors = validate(_deck(spec), allow_sample_content=True)
        assert not errors, "\n".join(errors)
        slide = _slide()
        s_split(slide, spec, 1)
        assert len(slide.shapes) > 5

    invalid = deepcopy(diagram_cards)
    invalid["left"] = {
        "type": "split", "heading": "入れ子",
        "left": {}, "right": {},
    }
    errors = validate(_deck(invalid))
    assert any("半幅未対応" in error for error in errors), errors

    unknown = deepcopy(diagram_cards)
    unknown["right"]["x"] = 1.0
    errors = validate(_deck(unknown))
    assert any("right.x" in error for error in errors), errors

    too_many_columns = deepcopy(diagram_cards)
    too_many_columns["left"]["diagram"]["cols"].append("extra")
    errors = validate(_deck(too_many_columns))
    assert any("半幅では2〜3件" in error for error in errors), errors

    dense = deepcopy(bullets_cards)
    dense["left"]["bullets"] = [
        ["半幅領域へ収まらない長い説明を繰り返して明示停止を確認する。" * 8, None]
        for _ in range(4)
    ]
    _must_fit_fail(dense)
    print("split layout tests passed")


if __name__ == "__main__":
    main()
